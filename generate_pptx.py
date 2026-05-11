import pptx
import pptx.util
import pptx.enum.shapes
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

# ── カラー定義 ──
BLACK   = RGBColor(0x0A, 0x0A, 0x0A)
BLACK2  = RGBColor(0x11, 0x11, 0x11)
BLACK3  = RGBColor(0x1A, 0x1A, 0x1A)
WHITE   = RGBColor(0xF5, 0xF5, 0xF0)
WHITE2  = RGBColor(0xCC, 0xCC, 0xCC)
GOLD    = RGBColor(0xC9, 0xA2, 0x27)
GOLD2   = RGBColor(0xE8, 0xC3, 0x5A)
GREEN   = RGBColor(0x1B, 0x43, 0x32)
GREEN2  = RGBColor(0x2D, 0x6A, 0x4F)
GREEN3  = RGBColor(0x52, 0xB7, 0x88)
RED2    = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

blank_layout = prs.slide_layouts[6]  # 完全空白

# ── ヘルパー関数 ──
def bg(slide, color):
    """スライド背景色を設定"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color=None, border_color=None, border_pt=1):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE rectangle
        int(x), int(y), int(w), int(h)
    )
    shape.line.color.rgb = border_color if border_color else (fill_color or BLACK)
    if border_color is None:
        shape.line.fill.background()  # no border
    else:
        shape.line.width = Pt(border_pt)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def txbox(slide, text, x, y, w, h,
          font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          font_name="Hiragino Sans", wrap=True):
    from pptx.util import Emu
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def txbox_multi(slide, lines, x, y, w, h,
                font_size=16, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, line_spacing=None,
                font_name="Hiragino Sans"):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line_text, fs, fc, fb in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(fs if fs else font_size)
        run.font.bold = fb if fb is not None else bold
        run.font.color.rgb = fc if fc else color
        run.font.name = font_name
    return tb

def gold_line(slide, x, y1, y2, width_pt=2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        int(x), int(y1), int(x), int(y2)
    )
    line.line.color.rgb = GOLD
    line.line.width = Pt(width_pt)

def hline(slide, x1, x2, y, color=GOLD, width_pt=1):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        int(x1), int(y), int(x2), int(y)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)

# ════════════════════════════════════════════════════════
# SLIDE 1 — 表紙
# ════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
bg(s1, BLACK)

# 左の金縦ライン
gold_line(s1, Inches(0.65), 0, H)

# 右側深緑グラデ代わりの矩形
rect(s1, W - Inches(2.8), 0, Inches(2.8), H, fill_color=RGBColor(0x0A, 0x1F, 0x15))

# PPPJ ロゴ部分
txbox(s1, "PEOPLE POWER POTENTIAL JAPAN", W - Inches(2.6), Inches(0.3), Inches(2.4), Inches(0.3),
      font_size=7, color=GOLD, align=PP_ALIGN.RIGHT)
txbox(s1, "PPPJ", W - Inches(1.8), Inches(0.55), Inches(1.6), Inches(0.6),
      font_size=32, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# eyebrow
txbox(s1, "HALAL × INBOUND SEMINAR 2026", Inches(1.0), Inches(1.2), Inches(8), Inches(0.4),
      font_size=11, color=GOLD)

# メインタイトル
txbox_multi(s1, [
    ('"選ばれる飲食店"は、', 40, WHITE, True),
    ('なぜ世界基準に', 48, WHITE, True),
    ('変わり始めているのか', 48, GOLD, True),
], Inches(1.0), Inches(1.8), Inches(8.5), Inches(2.4), align=PP_ALIGN.LEFT)

# サブコピー
rect(s1, Inches(1.0), Inches(4.3), Inches(3), Inches(0.04), fill_color=GOLD)
txbox(s1, '"美味しい"だけでは、選ばれなくなる時代へ。',
      Inches(1.0), Inches(4.5), Inches(8), Inches(0.45),
      font_size=16, color=WHITE2)

# 下部 stats（3つ）
stat_data = [
    ("3,687万人", "2024年 訪日外国人数\n（過去最大）"),
    ("20億人",    "世界ムスリム人口\n2030年：22億人"),
    ("1%以下",   "日本のハラル\n対応店舗割合"),
]
for i, (n, l) in enumerate(stat_data):
    sx = Inches(1.0) + i * Inches(3.0)
    hline(s1, sx, sx + Inches(2.4), Inches(5.6))
    txbox(s1, n, sx, Inches(5.7), Inches(2.4), Inches(0.5),
          font_size=24, bold=True, color=GOLD)
    txbox(s1, l, sx, Inches(6.2), Inches(2.4), Inches(0.7),
          font_size=10, color=WHITE2)

# ════════════════════════════════════════════════════════
# SLIDE 2 — 問題提起
# ════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
bg(s2, BLACK2)

# 左パネル（深緑）
rect(s2, 0, 0, Inches(4.6), H, fill_color=GREEN)

# PROBLEM タグ
rect(s2, Inches(0.5), Inches(0.7), Inches(1.3), Inches(0.3), fill_color=None, border_color=GOLD)
txbox(s2, "PROBLEM", Inches(0.52), Inches(0.72), Inches(1.25), Inches(0.28),
      font_size=9, color=GOLD, align=PP_ALIGN.CENTER)

txbox_multi(s2, [
    ('"食べたいのに、', 32, WHITE, True),
    ('食べられない"', 40, GOLD, True),
    ('外国人観光客が', 28, WHITE, True),
    ('増えています。', 28, WHITE, True),
], Inches(0.5), Inches(1.2), Inches(3.8), Inches(3.5))

txbox(s2, "日本食への期待は世界最高水準。\nしかし、その期待に応えられていない現実があります。",
      Inches(0.5), Inches(4.9), Inches(3.8), Inches(1.2),
      font_size=12, color=WHITE2)

# 右パネル
txbox(s2, "KEY DATA", Inches(5.0), Inches(0.7), Inches(1.5), Inches(0.3),
      font_size=9, color=GOLD)
hline(s2, Inches(5.0), Inches(5.9), Inches(1.1), color=GOLD)

cards = [
    ("78%",    "訪日外国人のうち「日本食を楽しむこと」を\n旅行目的に挙げた割合"),
    ("67%",    "ムスリム訪日客が「食事の制約」を\n最大の課題と回答した割合"),
    ("150万人", "年間訪日ムスリム観光客数（推計）\n2004年比：約7倍に増加"),
    ("22億人",  "2030年の世界ムスリム人口予測\n（世界人口の4人に1人）"),
]
cols = [0, 1, 0, 1]
rows = [0, 0, 1, 1]
for i, (n, l) in enumerate(cards):
    cx = Inches(5.0) + cols[i] * Inches(3.8)
    cy = Inches(1.4) + rows[i] * Inches(2.7)
    border_c = GOLD if i == 1 else RGBColor(0x40, 0x36, 0x12)
    fill_c   = RGBColor(0x1A, 0x16, 0x05) if i == 1 else RGBColor(0x18, 0x18, 0x18)
    rect(s2, cx, cy, Inches(3.4), Inches(2.3), fill_color=fill_c, border_color=border_c)
    txbox(s2, n, cx + Inches(0.2), cy + Inches(0.15), Inches(3.0), Inches(0.9),
          font_size=40, bold=True, color=GOLD)
    txbox(s2, l, cx + Inches(0.2), cy + Inches(1.0), Inches(3.0), Inches(1.1),
          font_size=11, color=WHITE2)

# ════════════════════════════════════════════════════════
# SLIDE 3 — 市場変化
# ════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
bg(s3, BLACK)

txbox(s3, "MARKET SHIFT", Inches(0.5), Inches(0.4), Inches(2.5), Inches(0.3),
      font_size=9, color=GOLD)

txbox_multi(s3, [
    ('飲食店の競争相手は、', 36, WHITE, True),
    ('"隣の店"ではなくなっています。', 36, GOLD, True),
], Inches(0.5), Inches(0.8), Inches(12), Inches(1.3), align=PP_ALIGN.CENTER)

txbox(s3, "選ばれる理由が、「美味しさ」から「安心感」へとシフトしています。",
      Inches(0.5), Inches(2.1), Inches(12), Inches(0.4),
      font_size=14, color=WHITE2, align=PP_ALIGN.CENTER)

card_data = [
    ("🔍", "Googleで選ばれる",   "来日前に検索。対応店だけが候補になる。"),
    ("📱", "SNSで探される",      "口コミが世界に広がる。"),
    ("🌍", "安心感で比較される", "「ここなら安心」が決め手になる。"),
    ("🗺️", "ガイドブックに載る", "海外ガイドアプリ掲載が新客を呼ぶ。"),
    ("🤝", "翻訳対応で差がつく", "多言語メニューが選ばれ方を変える。"),
]
for i, (icon, title, desc) in enumerate(card_data):
    cx = Inches(0.4) + i * Inches(2.55)
    rect(s3, cx, Inches(2.8), Inches(2.3), Inches(2.8),
         fill_color=RGBColor(0x18, 0x18, 0x18), border_color=RGBColor(0x30, 0x30, 0x30))
    txbox(s3, icon,  cx + Inches(0.8), Inches(3.0), Inches(0.8), Inches(0.5), font_size=28)
    txbox(s3, title, cx + Inches(0.15), Inches(3.6), Inches(2.0), Inches(0.45),
          font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s3, desc,  cx + Inches(0.1), Inches(4.1), Inches(2.1), Inches(1.2),
          font_size=11, color=WHITE2, align=PP_ALIGN.CENTER)

# 下部メッセージ
rect(s3, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.7),
     fill_color=None, border_color=GOLD)
txbox(s3, "美味しいだけでは、もう勝てない。",
      Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.7),
      font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 4 — 危機感
# ════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
bg(s4, RGBColor(0x0D, 0x08, 0x08))

rect(s4, 0, 0, Inches(4.8), H, fill_color=RGBColor(0x2C, 0x0A, 0x0A))

txbox(s4, "COMMON MISTAKE", Inches(0.4), Inches(0.6), Inches(2.5), Inches(0.3),
      font_size=9, color=RED2)

txbox_multi(s4, [
    ('"認証を', 40, WHITE, True),
    ('取っただけ"', 40, RED2, True),
    ('では、', 36, WHITE, True),
    ('お客様は来ません。', 28, WHITE, True),
], Inches(0.4), Inches(1.1), Inches(4.0), Inches(3.8))

txbox(s4, "ハラル認証は「入場券」に過ぎません。\n多くの飲食店が、この段階で止まっています。",
      Inches(0.4), Inches(5.1), Inches(4.0), Inches(1.2),
      font_size=12, color=WHITE2)

# 右パネル
txbox(s4, "よくある失敗パターン", Inches(5.2), Inches(0.6), Inches(7.0), Inches(0.4),
      font_size=12, color=RED2)
hline(s4, Inches(5.2), Inches(12.8), Inches(1.1), color=RGBColor(0x60, 0x20, 0x20))

fails = [
    "認証取得で満足して、情報発信をしない",
    "SNS・Googleの掲載情報が未整備のまま",
    "スタッフがイスラム文化を理解していない",
    "海外向けメニューや多言語対応がない",
    "食材・調味料の見直しが不完全",
    "現場オペレーションが続かない",
]
for i, f in enumerate(fails):
    fy = Inches(1.3) + i * Inches(0.72)
    rect(s4, Inches(5.2), fy, Inches(7.6), Inches(0.58),
         fill_color=RGBColor(0x1A, 0x0A, 0x0A), border_color=None)
    rect(s4, Inches(5.2), fy, Inches(0.04), Inches(0.58),
         fill_color=RGBColor(0x80, 0x20, 0x20))
    txbox(s4, f"✕  {f}", Inches(5.35), fy + Inches(0.1), Inches(7.3), Inches(0.45),
          font_size=14, color=WHITE2)

rect(s4, Inches(5.2), Inches(5.8), Inches(7.6), Inches(0.8),
     fill_color=RGBColor(0x1A, 0x14, 0x02), border_color=GOLD)
txbox(s4, "「知られない・伝わらない・続かない」——この3つが揃うと、お客様は来ない。",
      Inches(5.35), Inches(5.85), Inches(7.3), Inches(0.7),
      font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 5 — 成功パターン
# ════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
bg(s5, BLACK2)

txbox(s5, "SUCCESS PATTERN", Inches(0.5), Inches(0.35), Inches(2.8), Inches(0.3),
      font_size=9, color=GOLD)
txbox_multi(s5, [
    ('選ばれる店は、"認証"ではなく', 28, WHITE, True),
    ('"体験"を設計しています。', 32, GOLD, True),
], Inches(0.5), Inches(0.7), Inches(9.0), Inches(1.1))

# データピル
for i, (n, l) in enumerate([("87%", "ハラル対応店のうち\n来店客増加と回答"),
                              ("330%", "PPPJ支援店舗の\n最大売上UP実績")]):
    px = Inches(10.0) + i * Inches(1.6)
    rect(s5, px, Inches(0.35), Inches(1.4), Inches(1.1),
         fill_color=RGBColor(0x1A, 0x16, 0x05), border_color=GOLD)
    txbox(s5, n, px, Inches(0.38), Inches(1.4), Inches(0.55),
          font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txbox(s5, l, px + Inches(0.05), Inches(0.88), Inches(1.3), Inches(0.5),
          font_size=8, color=WHITE2, align=PP_ALIGN.CENTER)

# 8要素グリッド
elements = [
    ("🌏", "世界観の設計",  "店舗の雰囲気で\n「ここは安心」を伝える"),
    ("🗣️", "多言語接客",    "英語・アラビア語対応で\n信頼感が生まれる"),
    ("📸", "SNS・写真設計", "来店者が発信したくなる\nフォトスポット"),
    ("📍", "Google/MEO",   "検索されたとき\n確実に選択肢に入る"),
    ("📖", "ストーリー発信", "食材・調理・こだわりを\n世界に伝える"),
    ("🥩", "食材の信頼性",  "認証食材で「安心」を\n食材レベルから証明"),
    ("⭐", "口コミ設計",    "TripAdvisor等で\n口コミが連鎖する仕組み"),
    ("🔄", "継続運用",      "数値を追い\n改善し続ける"),
]
for i, (icon, name, desc) in enumerate(elements):
    ex = Inches(0.4) + (i % 4) * Inches(3.2)
    ey = Inches(1.9) + (i // 4) * Inches(1.75)
    rect(s5, ex, ey, Inches(2.9), Inches(1.55),
         fill_color=RGBColor(0x16, 0x16, 0x16), border_color=RGBColor(0x30, 0x28, 0x0A))
    rect(s5, ex, ey + Inches(1.45), Inches(2.9), Inches(0.08), fill_color=GOLD)
    txbox(s5, icon, ex + Inches(0.1), ey + Inches(0.1), Inches(0.5), Inches(0.5), font_size=22)
    txbox(s5, name, ex + Inches(0.65), ey + Inches(0.12), Inches(2.1), Inches(0.42),
          font_size=13, bold=True, color=WHITE)
    txbox(s5, desc, ex + Inches(0.1), ey + Inches(0.65), Inches(2.7), Inches(0.75),
          font_size=10, color=WHITE2)

# 事例3つ
cases = [("＋35%", "那覇市 ハラル対応店\nイスラム圏観光客増加"),
         ("＋26%", "金沢市 外国人客数増加\n客単価+¥1,200"),
         ("3ヶ月", "最短の投資回収期間\n（PPPJ支援実績）")]
for i, (n, l) in enumerate(cases):
    cx = Inches(0.4) + i * Inches(4.3)
    rect(s5, cx, Inches(5.5), Inches(4.0), Inches(0.75),
         fill_color=RGBColor(0x0E, 0x22, 0x1B), border_color=RGBColor(0x30, 0x70, 0x50))
    txbox(s5, n, cx + Inches(0.2), Inches(5.55), Inches(1.2), Inches(0.65),
          font_size=28, bold=True, color=GREEN3)
    txbox(s5, l, cx + Inches(1.5), Inches(5.55), Inches(2.3), Inches(0.65),
          font_size=11, color=WHITE2)

# ════════════════════════════════════════════════════════
# SLIDE 6 — ハラル市場の可能性
# ════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
bg(s6, BLACK)

# 左テキスト
txbox(s6, "MARKET OPPORTUNITY", Inches(0.5), Inches(0.4), Inches(3.5), Inches(0.3),
      font_size=9, color=GOLD)
txbox_multi(s6, [
    ('ハラル対応は、', 30, WHITE, True),
    ('"宗教対応"だけでは', 30, GOLD, True),
    ('ありません。', 30, WHITE, True),
], Inches(0.5), Inches(0.8), Inches(5.5), Inches(1.8))

rect(s6, Inches(0.5), Inches(2.8), Inches(5.3), Inches(0.03), fill_color=GOLD)
txbox_multi(s6, [
    ('"世界中のお客様に', 20, GOLD, True),
    ('選ばれる入口"になっています。', 20, GOLD, True),
], Inches(0.5), Inches(2.9), Inches(5.3), Inches(1.0))
rect(s6, Inches(0.5), Inches(3.9), Inches(5.3), Inches(0.03), fill_color=GOLD)

txbox(s6, "ムスリム市場は「特別対応」ではありません。\n今、最も成長速度の高い顧客セグメントです。\n宗教色を超えた「世界標準のホスピタリティ」です。",
      Inches(0.5), Inches(4.1), Inches(5.3), Inches(1.4),
      font_size=13, color=WHITE2)

# 右パネル（深緑）
rect(s6, Inches(6.2), 0, Inches(7.1), H, fill_color=GREEN)

mstats = [
    ("22億人", "2030年の世界ムスリム人口予測\n世界人口の4人に1人がムスリムへ。\n東南アジア・中東からの訪日客も急増中。"),
    ("150万人", "年間訪日ムスリム観光客数\n2004年の15万人から約7倍に急拡大。\nインバウンド復活でさらに加速中。"),
    ("1%",     "日本のハラル対応店舗割合\n2019年時点でわずか242店。\n需要に対して供給が圧倒的に不足。"),
]
for i, (n, l) in enumerate(mstats):
    my = Inches(0.55) + i * Inches(1.9)
    rect(s6, Inches(6.4), my, Inches(6.7), Inches(1.65),
         fill_color=RGBColor(0x0A, 0x22, 0x15), border_color=RGBColor(0x40, 0x80, 0x50))
    txbox(s6, n, Inches(6.6), my + Inches(0.15), Inches(2.0), Inches(0.75),
          font_size=40, bold=True, color=GOLD)
    txbox(s6, l, Inches(8.7), my + Inches(0.15), Inches(4.2), Inches(1.3),
          font_size=11, color=WHITE2)

rect(s6, Inches(6.4), Inches(6.2), Inches(6.7), Inches(0.7),
     fill_color=RGBColor(0x1A, 0x14, 0x02), border_color=GOLD)
txbox(s6, "今すぐ動けば、99%の市場に先手を打てる。",
      Inches(6.4), Inches(6.2), Inches(6.7), Inches(0.7),
      font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 7 — PPPJ 登場
# ════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
bg(s7, BLACK)

# 背景の大文字 PPPJ
txbox(s7, "PPPJ", Inches(1.5), Inches(1.5), Inches(10), Inches(4),
      font_size=200, bold=True, color=RGBColor(0x18, 0x14, 0x05), align=PP_ALIGN.CENTER)

rect(s7, Inches(3.5), Inches(0.5), Inches(2.5), Inches(0.38),
     fill_color=RGBColor(0x25, 0x08, 0x05), border_color=RED2)
txbox(s7, "NOT A CERTIFICATION COMPANY", Inches(3.5), Inches(0.5), Inches(2.5), Inches(0.38),
      font_size=9, color=RED2, align=PP_ALIGN.CENTER)

txbox_multi(s7, [
    ('PPPJは、', 32, WHITE2, False),
    ('"認証取得会社"ではありません。', 32, RED2, True),
], Inches(1.0), Inches(1.1), Inches(11.3), Inches(1.0), align=PP_ALIGN.CENTER)

txbox(s7, "↓", Inches(6.0), Inches(2.15), Inches(1.3), Inches(0.7),
      font_size=36, color=GOLD, align=PP_ALIGN.CENTER)

txbox_multi(s7, [
    ('"売れる仕組みまで', 44, WHITE, True),
    ('作る会社です。"', 44, GOLD, True),
], Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.8), align=PP_ALIGN.CENTER)

info = [("2011年", "設立年"), ("約400名", "従業員数"), ("日本・タイ", "2拠点体制"), ("木原 徹", "代表取締役社長")]
for i, (n, l) in enumerate(info):
    ix = Inches(1.2) + i * Inches(2.8)
    hline(s7, ix, ix + Inches(2.2), Inches(6.3))
    txbox(s7, n, ix, Inches(6.4), Inches(2.2), Inches(0.5),
          font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txbox(s7, l, ix, Inches(6.9), Inches(2.2), Inches(0.35),
          font_size=10, color=WHITE2, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 8 — PPPJの強み（ハブ）
# ════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
bg(s8, BLACK2)

txbox_multi(s8, [
    ('6つの柱で、あなたの店を', 26, WHITE, True),
    ('"選ばれ続ける店"にする', 26, GOLD, True),
], Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8), align=PP_ALIGN.CENTER)

# 中央ハブ
rect(s8, Inches(5.35), Inches(2.0), Inches(2.6), Inches(2.6),
     fill_color=GOLD, border_color=GOLD)
txbox(s8, "PPPJ", Inches(5.35), Inches(2.5), Inches(2.6), Inches(0.8),
      font_size=36, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
txbox(s8, "People Power\nPotential Japan", Inches(5.35), Inches(3.25), Inches(2.6), Inches(0.7),
      font_size=10, color=BLACK, align=PP_ALIGN.CENTER)

# 左側3項目
left_items = [
    ("📋 ハラル認証取得支援", "申請〜審査〜取得まで完全サポート"),
    ("🥩 ハラル食材の安定供給", "天幸牛・調味料など信頼食材ネットワーク"),
    ("👨‍🍳 メニュー開発サポート", "日本食の魅力を活かしたハラル対応設計"),
]
for i, (t, d) in enumerate(left_items):
    ly = Inches(1.1) + i * Inches(1.7)
    rect(s8, Inches(0.3), ly, Inches(4.6), Inches(1.45),
         fill_color=RGBColor(0x16, 0x16, 0x16), border_color=RGBColor(0x35, 0x2A, 0x0E))
    txbox(s8, t, Inches(0.45), ly + Inches(0.1), Inches(4.2), Inches(0.5),
          font_size=14, bold=True, color=WHITE)
    txbox(s8, d, Inches(0.45), ly + Inches(0.6), Inches(4.2), Inches(0.7),
          font_size=11, color=WHITE2)

# 右側3項目
right_items = [
    ("📱 SNS・集客支援", "インフルエンサー施策・海外メディア連携"),
    ("📍 Google / MEO 整備", "検索で「選ばれる店」になる情報設計"),
    ("🎓 スタッフ教育・採用", "インバウンド接客・イスラム文化研修"),
]
for i, (t, d) in enumerate(right_items):
    ry = Inches(1.1) + i * Inches(1.7)
    rect(s8, Inches(8.4), ry, Inches(4.6), Inches(1.45),
         fill_color=RGBColor(0x16, 0x16, 0x16), border_color=RGBColor(0x35, 0x2A, 0x0E))
    txbox(s8, t, Inches(8.55), ry + Inches(0.1), Inches(4.2), Inches(0.5),
          font_size=14, bold=True, color=WHITE)
    txbox(s8, d, Inches(8.55), ry + Inches(0.6), Inches(4.2), Inches(0.7),
          font_size=11, color=WHITE2)

# 下部まとめ
rect(s8, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.75),
     fill_color=RGBColor(0x1A, 0x14, 0x02), border_color=GOLD)
txbox(s8, "認証取得から集客・リピートまで——すべてつながっています。6つが揃うから売上に直結する。",
      Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.75),
      font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 9 — 天幸牛
# ════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
bg(s9, BLACK)

rect(s9, 0, 0, Inches(5.0), H, fill_color=RGBColor(0x1A, 0x0E, 0x00))

rect(s9, Inches(0.4), Inches(0.6), Inches(2.2), Inches(0.34),
     fill_color=None, border_color=GOLD)
txbox(s9, "HALAL WAGYU BRAND", Inches(0.4), Inches(0.6), Inches(2.2), Inches(0.34),
      font_size=9, color=GOLD, align=PP_ALIGN.CENTER)

txbox_multi(s9, [
    ('"安心して扱える和牛"を、', 30, WHITE, True),
    ('世界へ。', 38, GOLD, True),
], Inches(0.4), Inches(1.1), Inches(4.2), Inches(1.5))

txbox(s9, "天 幸 牛（てんこうぎゅう）", Inches(0.4), Inches(2.7), Inches(4.2), Inches(0.45),
      font_size=18, bold=True, color=GOLD2)

hline(s9, Inches(0.4), Inches(4.6), Inches(3.25), color=GOLD)

txbox(s9, "自然豊かな北海道・九州の大地で、\n一頭一頭に愛情を注いで育てられた黒毛和牛。\n\n命をいただく瞬間にも、清浄さを貫く。\nハラールの教えに則ったダビーハ方式で、\nムスリムによる祈りとともに屠畜された牛肉。\n\nそれが、天幸牛。",
      Inches(0.4), Inches(3.4), Inches(4.2), Inches(3.0),
      font_size=12, color=WHITE2)

# 右パネル
txbox(s9, "天幸牛 ブランド基準 — TENKOU WAGYU STANDARDS",
      Inches(5.4), Inches(0.5), Inches(7.5), Inches(0.4),
      font_size=12, color=GOLD)
hline(s9, Inches(5.4), Inches(12.8), Inches(1.0), color=GOLD)

specs = [
    ("産地・品種",   "北海道・九州産\n黒毛和種 限定"),
    ("等級・規模",   "A4・A5のみ\n提携牧場1万頭から厳選"),
    ("ハラル対応",   "ダビーハ方式屠畜\nハラル認証食肉処理場"),
    ("品質管理",     "トレーサビリティ完全対応\n牧場→加工→販売"),
]
for i, (lbl, val) in enumerate(specs):
    sc = Inches(5.4) + (i % 2) * Inches(3.7)
    sr = Inches(1.15) + (i // 2) * Inches(1.8)
    rect(s9, sc, sr, Inches(3.3), Inches(1.55),
         fill_color=RGBColor(0x16, 0x16, 0x16), border_color=RGBColor(0x35, 0x2A, 0x0E))
    rect(s9, sc, sr, Inches(0.06), Inches(1.55), fill_color=GOLD)
    txbox(s9, lbl, sc + Inches(0.18), sr + Inches(0.1), Inches(2.9), Inches(0.35),
          font_size=10, color=WHITE2)
    txbox(s9, val, sc + Inches(0.18), sr + Inches(0.5), Inches(2.9), Inches(0.9),
          font_size=14, bold=True, color=WHITE)

rect(s9, Inches(5.4), Inches(5.0), Inches(7.4), Inches(1.1),
     fill_color=RGBColor(0x1A, 0x14, 0x02), border_color=GOLD)
txbox(s9, "「美味しい和牛を出したい」と「ムスリム客に安心してもらいたい」を、\n食材の最上流から両立させる——それが天幸牛の使命です。",
      Inches(5.55), Inches(5.1), Inches(7.1), Inches(0.9),
      font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 10 — 導入フロー
# ════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
bg(s10, BLACK2)

txbox_multi(s10, [
    ('導入から集客まで、', 30, GOLD, True),
    ('一気通貫で伴走します', 30, WHITE, True),
], Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9), align=PP_ALIGN.CENTER)

txbox(s10, "「何から始めればいいかわからない」を、ゼロから解決するのが私たちの役割です。",
      Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
      font_size=14, color=WHITE2, align=PP_ALIGN.CENTER)

steps = [
    ("STEP 01", "無料ヒアリング",  "店舗の現状・客層・\n目標をヒアリング"),
    ("STEP 02", "全体設計",        "認証・メニュー・\n集客の全体プラン"),
    ("STEP 03", "認証取得",        "申請〜取得まで\n専門家が完全サポート"),
    ("STEP 04", "集客展開",        "多言語SNS・Google・\nインフルエンサー施策"),
    ("STEP 05", "継続改善",        "データ・数値をもとに\nPDCAを継続"),
]
for i, (num, title, desc) in enumerate(steps):
    sx = Inches(0.35) + i * Inches(2.55)
    fill_c = RGBColor(0x0E, 0x22, 0x1B) if i % 2 == 0 else RGBColor(0x16, 0x16, 0x16)
    border_c = RGBColor(0x30, 0x70, 0x50) if i % 2 == 0 else RGBColor(0x30, 0x28, 0x0A)
    rect(s10, sx, Inches(1.9), Inches(2.3), Inches(3.5), fill_color=fill_c, border_color=border_c)
    txbox(s10, num, sx + Inches(0.1), Inches(2.0), Inches(2.1), Inches(0.35),
          font_size=10, color=GOLD)
    txbox(s10, title, sx + Inches(0.1), Inches(2.4), Inches(2.1), Inches(0.5),
          font_size=16, bold=True, color=WHITE)
    txbox(s10, desc, sx + Inches(0.1), Inches(3.0), Inches(2.1), Inches(1.2),
          font_size=12, color=WHITE2)
    if i < 4:
        txbox(s10, "→", sx + Inches(2.3), Inches(2.9), Inches(0.28), Inches(0.5),
              font_size=20, color=GOLD, align=PP_ALIGN.CENTER)

rect(s10, Inches(3.5), Inches(5.9), Inches(6.3), Inches(0.65),
     fill_color=GOLD, border_color=GOLD)
txbox(s10, "▶  最短3ヶ月で投資回収の実績あり",
      Inches(3.5), Inches(5.9), Inches(6.3), Inches(0.65),
      font_size=18, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 11 — 実績
# ════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
bg(s11, BLACK)

# 左パネル
rect(s11, 0, 0, W / 2 - Inches(0.1), H, fill_color=RGBColor(0x0D, 0x1A, 0x14))

txbox(s11, "肉の匠 将泰庵 神田店\nHALAL Wagyu Yakiniku SHOUTAIAN",
      Inches(0.5), Inches(0.5), Inches(5.5), Inches(0.8),
      font_size=12, color=WHITE2, align=PP_ALIGN.CENTER)
txbox(s11, "ハラル認証取得：2025年2月",
      Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.35),
      font_size=11, color=GOLD, align=PP_ALIGN.CENTER)

txbox(s11, "178", Inches(0.5), Inches(1.8), Inches(5.5), Inches(2.8),
      font_size=160, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s11, "% UP", Inches(0.5), Inches(4.5), Inches(5.5), Inches(0.7),
      font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

txbox(s11, "3月・4月にハラルインフルエンサーを導入し認知を拡大。\n以降、常に100%越えをキープし続け、最大178%を記録。",
      Inches(0.5), Inches(5.4), Inches(5.5), Inches(1.0),
      font_size=12, color=WHITE2, align=PP_ALIGN.CENTER)

# 右パネル
rect(s11, W / 2 + Inches(0.1), 0, W / 2 - Inches(0.1), H, fill_color=RGBColor(0x0D, 0x15, 0x20))

txbox(s11, "しゃぶしゃぶ 将泰庵 ヨドバシHD池袋ビル店\nHALAL Wagyu ShabuShabu SHOUTAIAN",
      Inches(6.9), Inches(0.5), Inches(5.9), Inches(0.8),
      font_size=12, color=WHITE2, align=PP_ALIGN.CENTER)
txbox(s11, "ハラル認証取得：2025年6月",
      Inches(6.9), Inches(1.35), Inches(5.9), Inches(0.35),
      font_size=11, color=GOLD, align=PP_ALIGN.CENTER)

txbox(s11, "330", Inches(6.9), Inches(1.8), Inches(5.9), Inches(2.8),
      font_size=160, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s11, "% UP", Inches(6.9), Inches(4.5), Inches(5.9), Inches(0.7),
      font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

txbox(s11, "神田店の施策をそのまま池袋に展開。\n右肩上がりの推移を描き、最大330%という驚異の数字を達成。",
      Inches(6.9), Inches(5.4), Inches(5.9), Inches(1.0),
      font_size=12, color=WHITE2, align=PP_ALIGN.CENTER)

# 中央の仕切り線
gold_line(s11, W / 2, Inches(0.5), H - Inches(0.5))

# ════════════════════════════════════════════════════════
# SLIDE 12 — クロージング
# ════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(blank_layout)
bg(s12, BLACK)

txbox(s12, "THE FUTURE OF JAPANESE RESTAURANTS",
      Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.35),
      font_size=10, color=GOLD, align=PP_ALIGN.CENTER)

txbox_multi(s12, [
    ('これからの飲食店は、', 36, WHITE, True),
    ('"美味しい"だけでは選ばれなくなる。', 36, GOLD, True),
], Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.5), align=PP_ALIGN.CENTER)

txbox_multi(s12, [
    ('"誰に"——どんなお客様に来てほしいのか。', 18, WHITE, False),
    ('"どう安心を届けるか"——その設計ができた店が、', 18, WHITE, False),
    ('世界から選ばれ続ける店になります。', 18, WHITE, False),
    ('', 10, WHITE, False),
    ('ムスリム対応は「特別な配慮」ではなく、', 16, WHITE2, False),
    ('世界標準のホスピタリティです。', 16, WHITE2, False),
], Inches(1.5), Inches(2.7), Inches(10.3), Inches(2.5), align=PP_ALIGN.CENTER)

rect(s12, Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.8),
     fill_color=RGBColor(0x1A, 0x14, 0x02), border_color=GOLD)
txbox(s12, 'PPPJは、“世界中のお客様に選ばれる店づくり”を支援しています。',
      Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.8),
      font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

stats12 = [("4,000万人", "2025年 訪日外国人予測"),
           ("最大330%", "PPPJ支援店舗 売上UP実績"),
           ("最短3ヶ月", "投資回収の実績")]
for i, (n, l) in enumerate(stats12):
    bx = Inches(1.2) + i * Inches(3.7)
    hline(s12, bx, bx + Inches(3.0), Inches(6.4))
    txbox(s12, n, bx, Inches(6.5), Inches(3.0), Inches(0.5),
          font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txbox(s12, l, bx, Inches(7.0), Inches(3.0), Inches(0.35),
          font_size=10, color=WHITE2, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 13 — CTA
# ════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(blank_layout)
bg(s13, BLACK)

# 左パネル（金）
rect(s13, 0, 0, Inches(5.2), H, fill_color=GOLD)

txbox(s13, "FREE CONSULTATION", Inches(0.4), Inches(0.6), Inches(4.4), Inches(0.35),
      font_size=10, color=RGBColor(0x50, 0x40, 0x00), align=PP_ALIGN.LEFT)

txbox(s13, "まずは、\n話してみてください。",
      Inches(0.4), Inches(1.1), Inches(4.4), Inches(1.6),
      font_size=36, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

txbox(s13, "「うちの店でできるのか？」\n「規模が小さくても大丈夫？」\n「何から始めればいい？」\n\nその答えは、無料相談でわかります。\nゼロからご一緒します。",
      Inches(0.4), Inches(3.0), Inches(4.4), Inches(2.2),
      font_size=14, color=RGBColor(0x30, 0x28, 0x00), align=PP_ALIGN.LEFT)

rect(s13, Inches(0.4), Inches(5.5), Inches(2.4), Inches(0.55),
     fill_color=BLACK, border_color=BLACK)
txbox(s13, "✓  初回相談 無料",
      Inches(0.4), Inches(5.5), Inches(2.4), Inches(0.55),
      font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# 右パネル
txbox(s13, "CONSULTATION MENU — ご相談内容",
      Inches(5.6), Inches(0.5), Inches(7.3), Inches(0.35),
      font_size=11, color=GOLD)
hline(s13, Inches(5.6), Inches(12.9), Inches(0.95), color=GOLD)

consults = [
    "▶  ハラル認証導入のご相談",
    "▶  天幸牛（ハラル和牛）お取り扱いのご相談",
    "▶  インバウンド・SNS集客のご相談",
    "▶  Google / MEO 整備のご相談",
]
for i, c in enumerate(consults):
    txbox(s13, c, Inches(5.6), Inches(1.15) + i * Inches(0.5), Inches(7.3), Inches(0.45),
          font_size=14, color=WHITE2)

# 電話
rect(s13, Inches(5.6), Inches(3.3), Inches(7.3), Inches(0.85),
     fill_color=RGBColor(0x16, 0x16, 0x16), border_color=RGBColor(0x30, 0x28, 0x0A))
txbox(s13, "📞  お電話でのお問い合わせ（平日10:00〜18:00）",
      Inches(5.75), Inches(3.33), Inches(7.0), Inches(0.35),
      font_size=10, color=WHITE2)
txbox(s13, "047-481-8555",
      Inches(5.75), Inches(3.6), Inches(7.0), Inches(0.5),
      font_size=26, bold=True, color=WHITE)

# メール
rect(s13, Inches(5.6), Inches(4.3), Inches(7.3), Inches(0.85),
     fill_color=RGBColor(0x16, 0x16, 0x16), border_color=RGBColor(0x30, 0x28, 0x0A))
txbox(s13, "✉️  メールでのお問い合わせ",
      Inches(5.75), Inches(4.33), Inches(7.0), Inches(0.35),
      font_size=10, color=WHITE2)
txbox(s13, "office@pppj.co.jp",
      Inches(5.75), Inches(4.6), Inches(7.0), Inches(0.5),
      font_size=22, bold=True, color=WHITE)

# QRエリア
rect(s13, Inches(5.6), Inches(5.3), Inches(7.3), Inches(1.5),
     fill_color=RGBColor(0x16, 0x16, 0x16), border_color=RGBColor(0x40, 0x36, 0x12))
rect(s13, Inches(5.8), Inches(5.45), Inches(1.1), Inches(1.1),
     fill_color=WHITE, border_color=WHITE)
txbox(s13, "QR\nCODE", Inches(5.8), Inches(5.7), Inches(1.1), Inches(0.65),
      font_size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
txbox(s13, "株式会社PPPJ 公式サイト",
      Inches(7.05), Inches(5.45), Inches(5.7), Inches(0.45),
      font_size=14, bold=True, color=WHITE)
txbox(s13, "QRコードからそのままLINEまたは\nフォームでお問い合わせいただけます。",
      Inches(7.05), Inches(5.9), Inches(5.7), Inches(0.75),
      font_size=12, color=WHITE2)

# ════════════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════════════
output_path = "/home/user/20260511_ccctest/PPPJ_ハラルセミナー資料.pptx"
prs.save(output_path)
print(f"✓ 保存完了: {output_path}")
