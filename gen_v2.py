#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPPJ ハラルセミナー PowerPoint Generator v2
26 slides with charts, photo placeholders, and bright design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Color Definitions ──
NAVY   = RGBColor(0x1B, 0x3A, 0x5C)
GOLD   = RGBColor(0xC9, 0xA2, 0x27)
GREEN  = RGBColor(0x2D, 0x6A, 0x4F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAM  = RGBColor(0xFD, 0xFB, 0xF4)
GRAY   = RGBColor(0xC8, 0xCB, 0xD0)
LGRAY  = RGBColor(0xF0, 0xF0, 0xF0)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
BLACK  = RGBColor(0x11, 0x11, 0x11)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
GOLD2  = RGBColor(0xE8, 0xC3, 0x5A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

blank_layout = prs.slide_layouts[6]

# ── Helper Functions ──
def bg(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color=None, border_color=None, border_pt=1.0, radius=False):
    shape = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, font_size=18, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, v_anchor=None, wrap=True, italic=False):
    from pptx.util import Pt as Pt2
    from pptx.enum.text import PP_ALIGN as PA
    txb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    txb.text_frame.word_wrap = wrap
    if v_anchor:
        txb.text_frame.vertical_anchor = v_anchor
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt2(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def photo_placeholder(slide, x, y, w, h, label="📷 写真をここに挿入"):
    """Gray box with white dashed inner border and centered label text."""
    # Outer gray fill
    shape = rect(slide, x, y, w, h, fill_color=GRAY, border_color=None)
    # Inner white border (simulate by adding a slightly smaller rect with no fill, white dashed border)
    margin = Inches(0.12)
    inner = slide.shapes.add_shape(1,
        int(x + margin), int(y + margin),
        int(w - margin*2), int(h - margin*2))
    inner.fill.background()
    inner.line.color.rgb = WHITE
    inner.line.width = Pt(1.5)
    inner.line.dash_style = 4  # dash
    # Label text
    txt(slide, label,
        x + margin, y + h/2 - Inches(0.25),
        w - margin*2, Inches(0.5),
        font_size=13, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
    return shape

def header_band(slide, title_text, subtitle_text=None):
    """Navy header band at top with white title text."""
    band_h = Inches(1.3)
    rect(slide, 0, 0, W, band_h, fill_color=NAVY)
    # Gold accent line at bottom of header
    rect(slide, 0, int(band_h - Pt(3)), W, int(Pt(3)), fill_color=GOLD)
    # Title
    txt(slide, title_text,
        Inches(0.4), Inches(0.15),
        W - Inches(0.8), Inches(0.75),
        font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        txt(slide, subtitle_text,
            Inches(0.4), Inches(0.85),
            W - Inches(0.8), Inches(0.4),
            font_size=13, color=GOLD2, align=PP_ALIGN.LEFT)

def slide_number(slide, n):
    txt(slide, str(n),
        W - Inches(0.5), H - Inches(0.35),
        Inches(0.35), Inches(0.3),
        font_size=10, color=DGRAY, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, fill=WHITE, border=LGRAY, border_pt=1.0):
    return rect(slide, x, y, w, h, fill_color=fill, border_color=border, border_pt=border_pt)

def add_chart_bar(slide, x, y, w, h, categories, values, title=None, navy_bars=True):
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series("値", values)
    from pptx.util import Inches as In
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        int(x), int(y), int(w), int(h),
        chart_data
    ).chart
    chart.has_legend = False
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    # Bar color navy
    if navy_bars:
        series = chart.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = NAVY
        # Highlight max value bar in gold
        max_val = max(values)
        for i, pt in enumerate(series.points):
            if values[i] == max_val:
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = GOLD
    return chart

def add_chart_line(slide, x, y, w, h, categories, values, title=None):
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series("訪日外客数(万人)", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h),
        chart_data
    ).chart
    chart.has_legend = False
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    series = chart.series[0]
    series.format.line.color.rgb = NAVY
    series.format.line.width = Pt(2.5)
    return chart

# ═══════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)

# Photo placeholder full width top (3.8")
ph_h = Inches(3.8)
photo_placeholder(slide, 0, 0, W, ph_h, "📷 写真をここに挿入")

# Navy overlay strip at very top for logo area
rect(slide, 0, 0, W, Inches(0.55), fill_color=NAVY)
txt(slide, "PPPJ", Inches(0.25), Inches(0.08), Inches(1.5), Inches(0.4),
    font_size=20, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

# Gold accent line below photo
rect(slide, 0, int(ph_h), W, int(Pt(5)), fill_color=GOLD)

# Main title
title_y = ph_h + Inches(0.15)
txt(slide,
    "\"選ばれる飲食店\"は、なぜ世界基準に変わり始めているのか",
    Inches(0.5), title_y, W - Inches(1.0), Inches(1.2),
    font_size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Subtitle
sub_y = title_y + Inches(1.25)
txt(slide,
    "訪日外国人4,000万人時代を生き残るための「ハラル対応」戦略セミナー",
    Inches(0.5), sub_y, W - Inches(1.0), Inches(0.55),
    font_size=15, color=DGRAY, align=PP_ALIGN.LEFT)

# PPPJ badge
badge_x = W - Inches(2.8)
badge_y = sub_y + Inches(0.65)
rect(slide, badge_x, badge_y, Inches(2.5), Inches(0.9),
     fill_color=NAVY, border_color=GOLD, border_pt=2)
txt(slide, "PPPJ　プレミアム・プロ・ピープル・ジャパン",
    badge_x, badge_y + Inches(0.07), Inches(2.5), Inches(0.75),
    font_size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

slide_number(slide, 1)

# ═══════════════════════════════════════════════════
# SLIDE 2 — Agenda / Goal
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "本日のゴール", "この15分で、あなたの行動が変わります")

# 3-step flow
flow_y = Inches(1.6)
flow_h = Inches(1.4)
steps = [
    ("①", "現状を知る", "訪日客の波と\nハラル市場の実態"),
    ("②", "可能性を見る", "実績店舗の\n売上データを確認"),
    ("③", "行動を決める", "PPPJと共に\n最初の一歩を踏み出す"),
]
step_w = Inches(3.3)
gap = Inches(0.35)
total_w = step_w * 3 + gap * 2
start_x = (W - total_w) / 2

for i, (num, title_s, desc) in enumerate(steps):
    sx = start_x + i * (step_w + gap)
    # Card background
    card(slide, sx, flow_y, step_w, flow_h, fill=WHITE, border=NAVY, border_pt=2)
    # Top accent
    rect(slide, sx, flow_y, step_w, Inches(0.12), fill_color=GOLD)
    # Number
    txt(slide, num, sx, flow_y + Inches(0.15), step_w, Inches(0.45),
        font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # Title
    txt(slide, title_s, sx, flow_y + Inches(0.55), step_w, Inches(0.4),
        font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Desc
    txt(slide, desc, sx, flow_y + Inches(0.9), step_w, Inches(0.5),
        font_size=11, color=DGRAY, align=PP_ALIGN.CENTER)
    # Arrow between steps
    if i < 2:
        arr_x = sx + step_w + Inches(0.03)
        txt(slide, "→", arr_x, flow_y + Inches(0.5), gap - Inches(0.05), Inches(0.5),
            font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Big message
msg_y = flow_y + flow_h + Inches(0.5)
rect(slide, Inches(1.0), msg_y, W - Inches(2.0), Inches(1.5),
     fill_color=NAVY, border_color=GOLD, border_pt=2)
txt(slide,
    "今日15分で、売上が変わるきっかけをつかんでいってください。",
    Inches(1.2), msg_y + Inches(0.3),
    W - Inches(2.4), Inches(0.9),
    font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(slide, 2)

# ═══════════════════════════════════════════════════
# SLIDE 3 — Opening Question
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "あなたの店に、チャンスは来ていますか？")

content_y = Inches(1.4)
content_h = H - content_y - Inches(0.5)

# Photo placeholder left half
photo_placeholder(slide, 0, content_y, W/2 - Inches(0.1), content_h)

# Right side content
rx = W/2 + Inches(0.1)
rw = W/2 - Inches(0.3)

# Question 1
txt(slide, "先月、あなたの店に",
    rx, content_y + Inches(0.2), rw, Inches(0.45),
    font_size=18, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
txt(slide, "外国人のお客様は何人来ましたか？",
    rx, content_y + Inches(0.6), rw, Inches(0.5),
    font_size=20, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Gold divider
rect(slide, rx, int(content_y + Inches(1.15)), int(rw * 0.6), int(Pt(3)), fill_color=GOLD)

# Question 2
txt(slide, "その数字、来年は何倍になっていますか？",
    rx, content_y + Inches(1.3), rw, Inches(0.7),
    font_size=18, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Data area
data_y = content_y + Inches(2.2)
card(slide, rx, data_y, rw, Inches(1.8), fill=LGRAY, border=NAVY, border_pt=1)
txt(slide, "訪日外国人数の推移",
    rx + Inches(0.15), data_y + Inches(0.1), rw - Inches(0.3), Inches(0.4),
    font_size=12, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
txt(slide, "3,687万人",
    rx + Inches(0.15), data_y + Inches(0.4), rw * 0.5, Inches(0.5),
    font_size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
txt(slide, "（2024年実績）",
    rx + Inches(0.15), data_y + Inches(0.85), rw * 0.5, Inches(0.35),
    font_size=13, color=DGRAY, align=PP_ALIGN.LEFT)
txt(slide, "→",
    rx + rw * 0.5, data_y + Inches(0.45), Inches(0.5), Inches(0.4),
    font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(slide, "4,000万人",
    rx + rw * 0.5 + Inches(0.45), data_y + Inches(0.4), rw * 0.45, Inches(0.5),
    font_size=22, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
txt(slide, "（2025年予測）",
    rx + rw * 0.5 + Inches(0.45), data_y + Inches(0.85), rw * 0.45, Inches(0.35),
    font_size=13, color=DGRAY, align=PP_ALIGN.LEFT)

slide_number(slide, 3)

# ═══════════════════════════════════════════════════
# SLIDE 4 — Visitor Data with LINE CHART
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "訪日外国人数の推移（2015-2024）", "波は来ている。問題は、受け取れているかどうかだ。")

# Line chart left 2/3
chart_x = Inches(0.3)
chart_y = Inches(1.4)
chart_w = W * 0.62
chart_h = H - Inches(2.2)

years = ["2015","2016","2017","2018","2019","2020","2021","2022","2023","2024"]
values = (1974, 2404, 2869, 3119, 3188, 412, 25, 383, 2507, 3687)

add_chart_line(slide, chart_x, chart_y, chart_w, chart_h, years, values,
               title="訪日外客数（万人）")

# COVID note
txt(slide, "※2020-2022はコロナ禍による激減",
    chart_x, chart_y + chart_h + Inches(0.05), chart_w, Inches(0.3),
    font_size=10, color=DGRAY, italic=True, align=PP_ALIGN.LEFT)

# Photo placeholder right
ph_x = W * 0.64 + Inches(0.1)
ph_w = W - ph_x - Inches(0.2)
photo_placeholder(slide, ph_x, chart_y, ph_w, chart_h * 0.55)

# Key message box
msg_y = chart_y + chart_h * 0.6
rect(slide, ph_x, msg_y, ph_w, chart_h * 0.4,
     fill_color=NAVY, border_color=GOLD, border_pt=2)
txt(slide, "波は来ている。\n問題は、受け取れているかどうかだ。",
    ph_x + Inches(0.1), msg_y + Inches(0.15),
    ph_w - Inches(0.2), chart_h * 0.4 - Inches(0.3),
    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(slide, 4)

# ═══════════════════════════════════════════════════
# SLIDE 5 — 日本食
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "訪日外国人が求めるもの —「日本食」")

content_y = Inches(1.4)
half_w = W / 2

# Photo placeholder left
photo_placeholder(slide, 0, content_y, half_w - Inches(0.15), H - content_y - Inches(0.4))

# Right side
rx = half_w + Inches(0.1)
rw = W - rx - Inches(0.2)

# Big % number
txt(slide, "78%",
    rx, content_y + Inches(0.3), rw, Inches(1.5),
    font_size=80, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Underline
rect(slide, rx + rw*0.1, int(content_y + Inches(1.75)), int(rw * 0.8), int(Pt(4)), fill_color=GOLD)

txt(slide, "訪日外国人が",
    rx, content_y + Inches(1.9), rw, Inches(0.5),
    font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txt(slide, "「日本食を楽しむこと」を目的に来日",
    rx, content_y + Inches(2.35), rw, Inches(0.5),
    font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Source note
txt(slide, "出典：観光庁 訪日外国人消費動向調査",
    rx, content_y + Inches(3.1), rw, Inches(0.3),
    font_size=10, color=DGRAY, italic=True, align=PP_ALIGN.CENTER)

# Gold accent card
card_y = content_y + Inches(3.5)
rect(slide, rx, card_y, rw, Inches(0.8), fill_color=GOLD, border_color=None)
txt(slide, "日本食は最大の観光資源 —\nあなたの店はその恩恵を受けていますか？",
    rx + Inches(0.1), card_y + Inches(0.05), rw - Inches(0.2), Inches(0.7),
    font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_number(slide, 5)

# ═══════════════════════════════════════════════════
# SLIDE 6 — ムスリム問題
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "見えない壁 ——「ムスリムのお客様」が食べられない現実")

content_y = Inches(1.4)

# Photo placeholder top-right
ph_w = Inches(3.5)
ph_h = Inches(2.0)
photo_placeholder(slide, W - ph_w - Inches(0.2), content_y, ph_w, ph_h)

# Quote
quote_x = Inches(0.4)
quote_w = W - ph_w - Inches(0.9)
rect(slide, quote_x - Inches(0.05), content_y, Inches(0.08), Inches(1.5), fill_color=GOLD)
txt(slide, "\"食べたいのに、食べられない。\"",
    quote_x + Inches(0.1), content_y + Inches(0.2), quote_w, Inches(0.8),
    font_size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT, italic=True)
txt(slide, "——ムスリム旅行者が日本で直面する最大の課題",
    quote_x + Inches(0.1), content_y + Inches(0.95), quote_w, Inches(0.4),
    font_size=14, color=DGRAY, align=PP_ALIGN.LEFT)

# 3 data cards
card_y = content_y + Inches(2.3)
card_w = (W - Inches(0.8)) / 3 - Inches(0.2)
data_items = [
    ("67%", "訪日ムスリムが食事\nに困ったと回答"),
    ("150万人", "年間訪日ムスリム数\n（推計）"),
    ("22億人", "世界のムスリム人口\n世界人口の約28%"),
]
for i, (val, desc) in enumerate(data_items):
    cx = Inches(0.4) + i * (card_w + Inches(0.2))
    card(slide, cx, card_y, card_w, Inches(2.2), fill=WHITE, border=NAVY, border_pt=2)
    rect(slide, cx, card_y, card_w, Inches(0.1), fill_color=NAVY)
    txt(slide, val, cx, card_y + Inches(0.2), card_w, Inches(0.9),
        font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(slide, desc, cx, card_y + Inches(1.05), card_w, Inches(0.9),
        font_size=13, color=NAVY, align=PP_ALIGN.CENTER)

slide_number(slide, 6)

# ═══════════════════════════════════════════════════
# SLIDE 7 — 競争構造変化
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "競争構造の変化 —— ハラル対応が「差別化」から「必須」へ")

content_y = Inches(1.4)
col_w = (W - Inches(0.8)) / 2 - Inches(0.15)

# Before column
bx = Inches(0.4)
rect(slide, bx, content_y, col_w, Inches(0.55), fill_color=DGRAY)
txt(slide, "〈 BEFORE 〉ハラル対応なし",
    bx, content_y + Inches(0.05), col_w, Inches(0.45),
    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
before_items = [
    "✕  外国人グループから除外される",
    "✕  SNSでネガティブ投稿のリスク",
    "✕  予約・問い合わせゼロ",
    "✕  インバウンド売上ゼロ",
    "✕  国内客減少を補う手段なし",
]
for j, item in enumerate(before_items):
    by = content_y + Inches(0.6) + j * Inches(0.6)
    c = LGRAY if j % 2 == 0 else WHITE
    card(slide, bx, by, col_w, Inches(0.55), fill=c, border=GRAY, border_pt=0.5)
    txt(slide, item, bx + Inches(0.1), by + Inches(0.1), col_w - Inches(0.2), Inches(0.38),
        font_size=13, color=RED, align=PP_ALIGN.LEFT)

# After column
ax = bx + col_w + Inches(0.3)
rect(slide, ax, content_y, col_w, Inches(0.55), fill_color=NAVY)
txt(slide, "〈 AFTER 〉ハラル対応あり",
    ax, content_y + Inches(0.05), col_w, Inches(0.45),
    font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
after_items = [
    "✓  ムスリム客グループが来店",
    "✓  SNSで口コミ・拡散が生まれる",
    "✓  海外からの予約が増加",
    "✓  月商+100%以上の事例あり",
    "✓  先行者として市場を独占",
]
for j, item in enumerate(after_items):
    ay = content_y + Inches(0.6) + j * Inches(0.6)
    c = LGRAY if j % 2 == 0 else WHITE
    card(slide, ax, ay, col_w, Inches(0.55), fill=c, border=GRAY, border_pt=0.5)
    txt(slide, item, ax + Inches(0.1), ay + Inches(0.1), col_w - Inches(0.2), Inches(0.38),
        font_size=13, color=GREEN, align=PP_ALIGN.LEFT, bold=True)

# Smartphone mockup box (center)
sm_x = bx + col_w + Inches(0.03)
sm_y = content_y + Inches(1.5)
sm_w = Inches(0.22)
sm_h = Inches(1.2)
rect(slide, sm_x, sm_y, sm_w, sm_h, fill_color=NAVY, border_color=GOLD, border_pt=1)
txt(slide, "📱", sm_x, sm_y + Inches(0.35), sm_w, Inches(0.5),
    font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(slide, 7)

# ═══════════════════════════════════════════════════
# SLIDE 8 — ハラル対応現状
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "日本のハラル対応飲食店の現状 —— まだ1%以下")

content_y = Inches(1.4)

# Large infographic area
# Big "1%以下" text
txt(slide, "1%以下",
    Inches(0.5), content_y + Inches(0.3), Inches(5.5), Inches(2.0),
    font_size=72, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

txt(slide, "日本の飲食店のうち\nハラル認証取得店の割合",
    Inches(0.5), content_y + Inches(2.3), Inches(5.5), Inches(0.9),
    font_size=18, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Right side stats
rx = Inches(6.5)
rw = W - rx - Inches(0.3)

# 99% non-compliant
rect(slide, rx, content_y + Inches(0.2), rw, Inches(1.1), fill_color=NAVY, border_color=GOLD, border_pt=2)
txt(slide, "99%が非対応",
    rx + Inches(0.1), content_y + Inches(0.3), rw - Inches(0.2), Inches(0.7),
    font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Person icons (simplified visual)
icon_y = content_y + Inches(1.5)
txt(slide, "👤" * 10,
    rx, icon_y, rw, Inches(0.4),
    font_size=18, color=DGRAY, align=PP_ALIGN.CENTER)
txt(slide, "👤" * 10,
    rx, icon_y + Inches(0.4), rw, Inches(0.4),
    font_size=18, color=DGRAY, align=PP_ALIGN.CENTER)
txt(slide, "👤" * 9 + "🌟",
    rx, icon_y + Inches(0.8), rw, Inches(0.4),
    font_size=18, color=DGRAY, align=PP_ALIGN.CENTER)
txt(slide, "（100店中、認証あり: 1店未満）",
    rx, icon_y + Inches(1.25), rw, Inches(0.35),
    font_size=11, color=DGRAY, italic=True, align=PP_ALIGN.CENTER)

# Data box
rect(slide, rx, icon_y + Inches(1.7), rw, Inches(0.75), fill_color=LGRAY, border_color=NAVY)
txt(slide, "認証取得店舗数: 242店（2019年時点）\n全国飲食店数の0.04%",
    rx + Inches(0.1), icon_y + Inches(1.75), rw - Inches(0.2), Inches(0.65),
    font_size=12, color=NAVY, align=PP_ALIGN.CENTER)

# Bottom message
rect(slide, Inches(0.3), H - Inches(1.2), W - Inches(0.6), Inches(0.85),
     fill_color=GOLD, border_color=None)
txt(slide, "これは「危機」ではなく「チャンス」 —— 先に動いた者が市場を独占する",
    Inches(0.5), H - Inches(1.15), W - Inches(1.0), Inches(0.75),
    font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_number(slide, 8)

# ═══════════════════════════════════════════════════
# SLIDE 9 — 失敗パターン
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "よくある失敗パターン —— 認証だけでは売れない")

content_y = Inches(1.4)

# Left: failure cycle diagram
cycle_cx = Inches(3.5)
cycle_cy = Inches(3.8)
cycle_r_x = Inches(2.2)
cycle_r_y = Inches(1.5)

# 4 steps in "circle" layout (arranged as rectangle)
cycle_steps = [
    ("認証取得", Inches(1.0), Inches(1.7)),
    ("発信なし", Inches(4.8), Inches(1.7)),
    ("誰にも知られない", Inches(4.8), Inches(3.8)),
    ("来店ゼロ", Inches(1.0), Inches(3.8)),
]
for i, (label, cx, cy) in enumerate(cycle_steps):
    rect(slide, cx, cy + content_y - Inches(1.4), Inches(1.8), Inches(0.75),
         fill_color=LGRAY, border_color=RED, border_pt=2)
    txt(slide, label,
        cx, cy + content_y - Inches(1.4) + Inches(0.12),
        Inches(1.8), Inches(0.5),
        font_size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Arrows between them (simple text arrows)
arrow_positions = [
    (Inches(2.85), content_y + Inches(0.55), "→"),
    (Inches(5.1), content_y + Inches(1.7), "↓"),
    (Inches(2.85), content_y + Inches(2.65), "←"),
    (Inches(0.7), content_y + Inches(1.7), "↑"),
]
for ax2, ay2, ar in arrow_positions:
    txt(slide, ar, ax2, ay2, Inches(0.5), Inches(0.5),
        font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txt(slide, "「認証だけ」\n失敗サイクル",
    Inches(2.0), content_y + Inches(1.5), Inches(2.5), Inches(0.8),
    font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Right: 6 failure items
rx = Inches(7.0)
rw = W - rx - Inches(0.3)
fail_items = [
    "✕  認証を取ったが告知しなかった",
    "✕  日本語のみのメニュー",
    "✕  SNS・Googleが未整備",
    "✕  スタッフが説明できない",
    "✕  ハラル食材の調達が不安定",
    "✕  外国語対応ゼロ",
]
txt(slide, "具体的な失敗例",
    rx, content_y + Inches(0.1), rw, Inches(0.45),
    font_size=15, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
rect(slide, rx, int(content_y + Inches(0.5)), int(rw * 0.6), int(Pt(3)), fill_color=GOLD)
for j, item in enumerate(fail_items):
    fy = content_y + Inches(0.65) + j * Inches(0.65)
    c = LGRAY if j % 2 == 0 else WHITE
    card(slide, rx, fy, rw, Inches(0.6), fill=c, border=GRAY, border_pt=0.5)
    txt(slide, item, rx + Inches(0.1), fy + Inches(0.1), rw - Inches(0.2), Inches(0.42),
        font_size=13, color=RED, align=PP_ALIGN.LEFT)

slide_number(slide, 9)

# ═══════════════════════════════════════════════════
# SLIDE 10 — Timeline Overview
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "PPPJ 支援店舗タイムライン", "2025年〜2026年 実績店舗の軌跡")

content_y = Inches(1.4)

# Horizontal timeline line
line_y = Inches(3.8)
rect(slide, Inches(0.5), line_y, W - Inches(1.0), Inches(0.08), fill_color=NAVY)

stores = [
    ("2025.02", "将泰庵\n神田店"),
    ("2025.06", "将泰庵\n池袋店"),
    ("2025.10", "銀座のステーキ\n恵比寿店"),
    ("2025.11", "銀座のステーキ\n新宿店"),
    ("2025.11", "銀座のステーキ\n渋谷店"),
    ("2026.01", "川品\n浅草店"),
]
n = len(stores)
spacing = (W - Inches(1.0)) / (n - 1)

for i, (date, name) in enumerate(stores):
    sx = Inches(0.5) + i * spacing
    # Dot on timeline
    dot_r = Inches(0.15)
    rect(slide, int(sx - dot_r/2), int(line_y - dot_r/2 + Inches(0.04)),
         int(dot_r), int(dot_r),
         fill_color=GOLD if i % 2 == 0 else NAVY, border_color=WHITE, border_pt=1)
    # Alternating above/below
    if i % 2 == 0:
        card_y = line_y - Inches(2.0)
        card(slide, int(sx - Inches(1.0)), int(card_y), Inches(2.0), Inches(1.7),
             fill=WHITE, border=NAVY, border_pt=2)
        txt(slide, date, int(sx - Inches(0.9)), int(card_y + Inches(0.1)),
            Inches(1.8), Inches(0.35), font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, name, int(sx - Inches(0.9)), int(card_y + Inches(0.45)),
            Inches(1.8), Inches(0.8), font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # Connector line
        rect(slide, int(sx - Inches(0.02)), int(card_y + Inches(1.7)),
             int(Inches(0.04)), int(Inches(0.3) + dot_r/2),
             fill_color=NAVY)
    else:
        card_y = line_y + Inches(0.4)
        rect(slide, int(sx - Inches(0.02)), int(line_y + dot_r/2 + Inches(0.04)),
             int(Inches(0.04)), int(Inches(0.3)),
             fill_color=NAVY)
        card(slide, int(sx - Inches(1.0)), int(card_y), Inches(2.0), Inches(1.7),
             fill=WHITE, border=GOLD, border_pt=2)
        txt(slide, date, int(sx - Inches(0.9)), int(card_y + Inches(0.1)),
            Inches(1.8), Inches(0.35), font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, name, int(sx - Inches(0.9)), int(card_y + Inches(0.45)),
            Inches(1.8), Inches(0.8), font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_number(slide, 10)

# ═══════════════════════════════════════════════════
# SLIDE 11 — 将泰庵神田 BAR CHART
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "【実績①】将泰庵 神田店 —— ハラル対応後の売上推移")

content_y = Inches(1.4)

# Photo placeholder top strip
photo_placeholder(slide, 0, content_y, W, Inches(0.9), "📷 店舗写真をここに挿入")

chart_y = content_y + Inches(0.95)
chart_w = W * 0.62
chart_h = H - chart_y - Inches(0.8)

months = ["1月","2月","3月","4月","5月","6月","7月","8月","9月","10月"]
values_sho = (100, 104, 163, 137, 129, 154, 138, 146, 178, 174)

add_chart_bar(slide, Inches(0.3), chart_y, chart_w, chart_h,
              months, values_sho, title="ハラル導入後 月別売上指数（1月=100）")

# Right side: big number + notes
rx = W * 0.64 + Inches(0.1)
rw = W - rx - Inches(0.2)

rect(slide, rx, chart_y, rw, Inches(1.3), fill_color=NAVY, border_color=GOLD, border_pt=3)
txt(slide, "178% UP",
    rx + Inches(0.1), chart_y + Inches(0.1), rw - Inches(0.2), Inches(0.8),
    font_size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(slide, "最高月（9月）実績",
    rx + Inches(0.1), chart_y + Inches(0.85), rw - Inches(0.2), Inches(0.35),
    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

note_y = chart_y + Inches(1.4)
card(slide, rx, note_y, rw, Inches(1.2), fill=LGRAY, border=NAVY)
txt(slide, "🌟 インフルエンサー施策\n3-4月にムスリム\nインフルエンサーを導入\n→ SNS拡散で来客急増",
    rx + Inches(0.1), note_y + Inches(0.1), rw - Inches(0.2), Inches(1.0),
    font_size=11, color=NAVY, align=PP_ALIGN.LEFT)

card(slide, rx, note_y + Inches(1.25), rw, Inches(0.9), fill=WHITE, border=GREEN)
txt(slide, "✓ 認証取得: 2025年2月\n✓ ハラル食材導入完了\n✓ 英語メニュー整備",
    rx + Inches(0.1), note_y + Inches(1.3), rw - Inches(0.2), Inches(0.8),
    font_size=11, color=GREEN, align=PP_ALIGN.LEFT)

slide_number(slide, 11)

# ═══════════════════════════════════════════════════
# SLIDE 12 — 神田 Success Factors
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "将泰庵 神田店 —— 成功の三段階メカニズム")

content_y = Inches(1.4)

# 3-box causation diagram left side
box_w = Inches(2.8)
box_h = Inches(2.2)
gap_x = Inches(0.3)
boxes = [
    ("①", "認証取得", "（2025年2月）",
     "・PPPJ認証取得\n・ハラル食材切替\n・スタッフ教育完了"),
    ("②", "インフルエンサー\n導入", "（3-4月）",
     "・ムスリム系\n　インフルエンサー招致\n・SNS拡散施策\n・Google MEO強化"),
    ("③", "常時100%越え\n定着", "（5月〜）",
     "・口コミ効果持続\n・リピーター獲得\n・グループ予約増加"),
]

total_box_w = box_w * 3 + gap_x * 2
start_bx = Inches(0.3)

for i, (num, title_b, sub, desc) in enumerate(boxes):
    bx = start_bx + i * (box_w + gap_x)
    by = content_y + Inches(0.3)
    card(slide, bx, by, box_w, box_h, fill=WHITE, border=NAVY, border_pt=2)
    # Top color
    c = [NAVY, GOLD, GREEN][i]
    rect(slide, bx, by, box_w, Inches(0.12), fill_color=c)
    txt(slide, num, bx, by + Inches(0.15), box_w, Inches(0.45),
        font_size=20, bold=True, color=c, align=PP_ALIGN.CENTER)
    txt(slide, title_b, bx, by + Inches(0.55), box_w, Inches(0.55),
        font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(slide, sub, bx, by + Inches(1.05), box_w, Inches(0.3),
        font_size=11, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    txt(slide, desc, bx + Inches(0.1), by + Inches(1.35), box_w - Inches(0.2), Inches(0.75),
        font_size=11, color=DGRAY, align=PP_ALIGN.LEFT)
    # Arrow
    if i < 2:
        arr_x = bx + box_w + Inches(0.03)
        txt(slide, "→", arr_x, by + box_h/2 - Inches(0.3), gap_x - Inches(0.05), Inches(0.5),
            font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Result bar below boxes
res_y = content_y + Inches(0.3) + box_h + Inches(0.25)
rect(slide, start_bx, res_y, total_box_w, Inches(0.65), fill_color=GOLD)
txt(slide, "→ 結果: 最大178%売上アップ ／ 月平均145%以上を維持",
    start_bx + Inches(0.2), res_y + Inches(0.1), total_box_w - Inches(0.4), Inches(0.45),
    font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Photo placeholder right
ph_x = start_bx + total_box_w + Inches(0.3)
ph_w = W - ph_x - Inches(0.2)
photo_placeholder(slide, ph_x, content_y, ph_w, H - content_y - Inches(0.5))

slide_number(slide, 12)

# ═══════════════════════════════════════════════════
# SLIDE 13 — 将泰庵池袋 BAR CHART
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "【実績②】将泰庵 池袋店 —— 神田モデルの横展開")

content_y = Inches(1.4)

# Photo placeholder top strip
photo_placeholder(slide, 0, content_y, W, Inches(0.9), "📷 店舗写真をここに挿入")

chart_y = content_y + Inches(0.95)
chart_w = W * 0.62
chart_h = H - chart_y - Inches(0.8)

months_ike = ["6月","7月","8月","9月","10月","11月"]
values_ike = (100, 233, 185, 185, 271, 330)

add_chart_bar(slide, Inches(0.3), chart_y, chart_w, chart_h,
              months_ike, values_ike, title="ハラル導入後 月別売上指数（6月=100）")

# Right side
rx = W * 0.64 + Inches(0.1)
rw = W - rx - Inches(0.2)

rect(slide, rx, chart_y, rw, Inches(1.3), fill_color=NAVY, border_color=GOLD, border_pt=3)
txt(slide, "330% UP",
    rx + Inches(0.1), chart_y + Inches(0.1), rw - Inches(0.2), Inches(0.8),
    font_size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(slide, "最高月（11月）実績",
    rx + Inches(0.1), chart_y + Inches(0.85), rw - Inches(0.2), Inches(0.35),
    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

note_y = chart_y + Inches(1.4)
card(slide, rx, note_y, rw, Inches(1.1), fill=LGRAY, border=GOLD)
txt(slide, "神田モデルを池袋に転用\n・認証: 2025年6月\n・即月233%を達成\n・わずか6ヶ月で330%に",
    rx + Inches(0.1), note_y + Inches(0.1), rw - Inches(0.2), Inches(0.9),
    font_size=12, color=NAVY, align=PP_ALIGN.LEFT)

card(slide, rx, note_y + Inches(1.2), rw, Inches(0.7), fill=WHITE, border=GREEN)
txt(slide, "✓ 再現性が証明された\n✓ PPPJの仕組みが機能",
    rx + Inches(0.1), note_y + Inches(1.25), rw - Inches(0.2), Inches(0.6),
    font_size=12, color=GREEN, align=PP_ALIGN.LEFT, bold=True)

slide_number(slide, 13)

# ═══════════════════════════════════════════════════
# SLIDE 14 — 池袋 Reproducibility
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "再現性こそがPPPJの強み —— モデルの転用")

content_y = Inches(1.4)

# Reproduction model diagram
# Left box: 神田モデル
bw = Inches(4.5)
bh = Inches(4.2)
bx1 = Inches(0.5)
by1 = content_y + Inches(0.2)

card(slide, bx1, by1, bw, bh, fill=WHITE, border=NAVY, border_pt=3)
rect(slide, bx1, by1, bw, Inches(0.6), fill_color=NAVY)
txt(slide, "神田モデル（原型）", bx1, by1 + Inches(0.1), bw, Inches(0.45),
    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
model_items = [
    "① ハラル認証取得",
    "② 食材をハラル対応に切替",
    "③ インフルエンサー施策",
    "④ SNS・Google MEO強化",
    "⑤ スタッフ教育・英語メニュー",
    "⑥ 継続的なデータ改善",
]
for j, item in enumerate(model_items):
    iy = by1 + Inches(0.75) + j * Inches(0.5)
    txt(slide, item, bx1 + Inches(0.2), iy, bw - Inches(0.4), Inches(0.42),
        font_size=13, color=NAVY, align=PP_ALIGN.LEFT,
        bold=(j == 0))

# Arrow
ax_m = bx1 + bw + Inches(0.1)
txt(slide, "→\n転用", ax_m, content_y + Inches(1.8), Inches(0.7), Inches(0.8),
    font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Right box: 池袋
bx2 = bx1 + bw + Inches(0.9)
card(slide, bx2, by1, bw, bh, fill=WHITE, border=GOLD, border_pt=3)
rect(slide, bx2, by1, bw, Inches(0.6), fill_color=GOLD)
txt(slide, "池袋転用（6ヶ月後）", bx2, by1 + Inches(0.1), bw, Inches(0.45),
    font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
result_items = [
    "→ 翌月: 233%達成",
    "→ 3ヶ月: 185%安定",
    "→ 5ヶ月: 271%",
    "→ 6ヶ月: 330%（最高）",
    "",
    "★ 同じ仕組みで同じ結果",
]
for j, item in enumerate(result_items):
    iy = by1 + Inches(0.75) + j * Inches(0.5)
    c = GREEN if item.startswith("→") else NAVY
    if item.startswith("★"):
        c = GOLD
    txt(slide, item, bx2 + Inches(0.2), iy, bw - Inches(0.4), Inches(0.42),
        font_size=13, color=c, align=PP_ALIGN.LEFT, bold=(j >= 5))

# Bottom message
rect(slide, Inches(0.3), H - Inches(1.0), W - Inches(0.6), Inches(0.65),
     fill_color=NAVY, border_color=GOLD, border_pt=2)
txt(slide, "再現性こそがPPPJの強み —— あなたの店でも同じ結果が出せます",
    Inches(0.5), H - Inches(0.95), W - Inches(1.0), Inches(0.55),
    font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

slide_number(slide, 14)

# ═══════════════════════════════════════════════════
# SLIDE 15 — 銀座のステーキ 3店舗
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "【実績③】銀座のステーキ 3店舗展開", "高単価×体験×ハラルの掛け算")

content_y = Inches(1.4)

# Photo placeholder top strip
photo_placeholder(slide, 0, content_y, W, Inches(1.0), "📷 店舗・料理写真をここに挿入")

store_y = content_y + Inches(1.1)
stores3 = [
    ("恵比寿店", "2025年10月オープン", "#1\n高級住宅街立地\n外国人富裕層が多い\nエリア"),
    ("新宿店", "2025年11月オープン", "#2\n最大集客エリア\n新宿駅周辺\nインバウンド最前線"),
    ("渋谷店", "2025年11月オープン", "#3\n若年富裕層向け\nトレンド発信地\nSNS拡散効果大"),
]
sw = (W - Inches(0.8)) / 3 - Inches(0.15)
for i, (name, date, desc) in enumerate(stores3):
    sx = Inches(0.4) + i * (sw + Inches(0.15))
    sy = store_y
    sh = H - sy - Inches(0.5)
    card(slide, sx, sy, sw, sh, fill=WHITE, border=NAVY, border_pt=2)
    rect(slide, sx, sy, sw, Inches(0.12), fill_color=GOLD)
    txt(slide, name, sx, sy + Inches(0.15), sw, Inches(0.5),
        font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(slide, date, sx, sy + Inches(0.65), sw, Inches(0.35),
        font_size=12, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    txt(slide, desc, sx + Inches(0.1), sy + Inches(1.05), sw - Inches(0.2), Inches(1.8),
        font_size=14, color=DGRAY, align=PP_ALIGN.CENTER)
    # Gold accent badge
    rect(slide, sx + sw/2 - Inches(0.8), sy + sh - Inches(0.75), Inches(1.6), Inches(0.5),
         fill_color=LGRAY, border_color=GOLD, border_pt=1)
    txt(slide, "ハラル認証済", sx + sw/2 - Inches(0.75), sy + sh - Inches(0.7),
        Inches(1.5), Inches(0.4), font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Bottom message
rect(slide, Inches(0.3), H - Inches(1.0), W - Inches(0.6), Inches(0.6),
     fill_color=NAVY, border_color=None)
txt(slide, "高単価×体験×ハラルの掛け算 —— プレミアムゾーンでの先行展開",
    Inches(0.5), H - Inches(0.95), W - Inches(1.0), Inches(0.5),
    font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

slide_number(slide, 15)

# ═══════════════════════════════════════════════════
# SLIDE 16 — 川品浅草
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "【実績④】川品 浅草店 —— 訪日外国人No.1エリアに出店")

content_y = Inches(1.4)

half_w = W / 2 - Inches(0.15)

# Left: 2 photo placeholders stacked
ph_h1 = (H - content_y - Inches(0.5)) / 2 - Inches(0.1)
photo_placeholder(slide, 0, content_y, half_w, ph_h1, "📷 浅草の街並み写真")
photo_placeholder(slide, 0, content_y + ph_h1 + Inches(0.1), half_w, ph_h1, "📷 料理写真")

# Right side
rx = half_w + Inches(0.3)
rw = W - rx - Inches(0.2)

# Header info card
card(slide, rx, content_y + Inches(0.1), rw, Inches(2.0), fill=WHITE, border=NAVY, border_pt=2)
rect(slide, rx, content_y + Inches(0.1), rw, Inches(0.12), fill_color=GOLD)
txt(slide, "川品 浅草店", rx + Inches(0.1), content_y + Inches(0.25), rw - Inches(0.2), Inches(0.55),
    font_size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
txt(slide, "オープン: 2026年1月\nジャンル: 和食・鍋・しゃぶしゃぶ\nエリア: 東京・浅草（台東区）\nPPPJ認証: 2025年12月取得",
    rx + Inches(0.1), content_y + Inches(0.82), rw - Inches(0.2), Inches(1.1),
    font_size=13, color=DGRAY, align=PP_ALIGN.LEFT)

# Key message
msg_y = content_y + Inches(2.25)
rect(slide, rx, msg_y, rw, Inches(1.0), fill_color=NAVY, border_color=GOLD, border_pt=2)
txt(slide, "訪日外国人来街数\n日本最多エリア・浅草",
    rx + Inches(0.1), msg_y + Inches(0.1), rw - Inches(0.2), Inches(0.8),
    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Stats
stat_y = msg_y + Inches(1.1)
card(slide, rx, stat_y, rw, Inches(0.75), fill=LGRAY, border=NAVY)
txt(slide, "年間来街者数: 約3,000万人以上\n外国人比率: 約60%（都内最高水準）",
    rx + Inches(0.1), stat_y + Inches(0.1), rw - Inches(0.2), Inches(0.6),
    font_size=13, color=NAVY, align=PP_ALIGN.LEFT)

card(slide, rx, stat_y + Inches(0.85), rw, Inches(0.75), fill=WHITE, border=GOLD)
txt(slide, "ムスリム旅行者に最も認知された\n観光地の一つ ——「浅草 ハラル」検索急増中",
    rx + Inches(0.1), stat_y + Inches(0.9), rw - Inches(0.2), Inches(0.65),
    font_size=12, color=NAVY, align=PP_ALIGN.LEFT)

slide_number(slide, 16)

# ═══════════════════════════════════════════════════
# SLIDE 17 — ROI まとめ
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "実績サマリー —— ROI 比較表", "動かないことが、最大のリスク。")

content_y = Inches(1.4)

# Table
table_x = Inches(0.4)
table_w = W * 0.65
row_h = Inches(0.75)
col_widths = [Inches(2.5), Inches(1.8), Inches(1.8), Inches(2.5)]
headers = ["店舗", "対応前", "最大実績", "ポイント"]

# Header row
hx = table_x
for k, (hdr, cw) in enumerate(zip(headers, col_widths)):
    rect(slide, int(hx), int(content_y), int(cw), int(row_h * 0.65),
         fill_color=NAVY, border_color=WHITE, border_pt=0.5)
    txt(slide, hdr, int(hx + Inches(0.05)), int(content_y + Inches(0.12)),
        int(cw - Inches(0.1)), int(row_h * 0.5),
        font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hx += cw

rows = [
    ("将泰庵 神田店", "基準:100", "178%", "インフルエンサー\n×認証で急拡大"),
    ("将泰庵 池袋店", "基準:100", "330%", "神田モデル転用\n再現性を証明"),
    ("銀座のステーキ\n恵比寿・新宿・渋谷", "新規出店", "展開中", "高単価プレミアム\n戦略で先行獲得"),
    ("川品 浅草店", "2026.01開始", "拡大中", "浅草No.1立地\n最多インバウンド"),
]

for r, row_data in enumerate(rows):
    ry = content_y + row_h * 0.65 + r * row_h
    rx2 = table_x
    fill = WHITE if r % 2 == 0 else LGRAY
    for k, (cell, cw) in enumerate(zip(row_data, col_widths)):
        rect(slide, int(rx2), int(ry), int(cw), int(row_h),
             fill_color=fill, border_color=GRAY, border_pt=0.5)
        fc = GOLD if (k == 2 and cell not in ["展開中","拡大中","新規出店","2026.01開始"]) else NAVY
        txt(slide, cell, int(rx2 + Inches(0.05)), int(ry + Inches(0.12)),
            int(cw - Inches(0.1)), int(row_h - Inches(0.1)),
            font_size=13, color=fc, align=PP_ALIGN.CENTER,
            bold=(k == 2))
        rx2 += cw

# Big message right side
mrx = table_x + sum(col_widths) + Inches(0.3)
mrw = W - mrx - Inches(0.2)
rect(slide, mrx, content_y, mrw, H - content_y - Inches(0.5),
     fill_color=NAVY, border_color=GOLD, border_pt=3)
txt(slide, "動かないことが、\n\n最大のリスク。",
    mrx + Inches(0.2), content_y + Inches(0.8), mrw - Inches(0.4), Inches(2.5),
    font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "競合が先に動いた瞬間\nあなたの市場は失われる",
    mrx + Inches(0.2), content_y + Inches(3.3), mrw - Inches(0.4), Inches(0.9),
    font_size=13, color=GOLD, align=PP_ALIGN.CENTER)

slide_number(slide, 17)

# ═══════════════════════════════════════════════════
# SLIDE 18 — PPPJとは
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "PPPJとは —— 私たちについて")

content_y = Inches(1.4)
half_w = W / 2 - Inches(0.1)

# Left: company info
lx = Inches(0.4)
lw = half_w - Inches(0.3)

txt(slide, "プレミアム・プロ・ピープル・ジャパン",
    lx, content_y + Inches(0.1), lw, Inches(0.5),
    font_size=15, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
txt(slide, "Premium Pro People Japan Co., Ltd.",
    lx, content_y + Inches(0.55), lw, Inches(0.4),
    font_size=12, color=DGRAY, italic=True, align=PP_ALIGN.LEFT)

rect(slide, lx, int(content_y + Inches(1.0)), int(lw * 0.4), int(Pt(3)), fill_color=GOLD)

company_info = [
    ("設立", "2011年"),
    ("従業員", "約400名"),
    ("拠点", "日本・タイの2拠点"),
    ("事業", "和牛総合商社"),
    ("認証実績", "飲食店ハラル対応支援"),
]
for j, (label, val) in enumerate(company_info):
    iy = content_y + Inches(1.15) + j * Inches(0.62)
    rect(slide, lx, int(iy), int(lw), int(Inches(0.58)),
         fill_color=LGRAY if j % 2 == 0 else WHITE, border_color=GRAY, border_pt=0.5)
    txt(slide, label + ":", lx + Inches(0.1), iy + Inches(0.1), Inches(1.2), Inches(0.4),
        font_size=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    txt(slide, val, lx + Inches(1.3), iy + Inches(0.1), lw - Inches(1.4), Inches(0.4),
        font_size=13, color=DGRAY, align=PP_ALIGN.LEFT)

# Founder message
msg_y = content_y + Inches(1.15) + 5 * Inches(0.62) + Inches(0.2)
rect(slide, lx - Inches(0.05), msg_y, Inches(0.08), Inches(1.2), fill_color=GOLD)
txt(slide, "「ムスリムの方々に安心して日本の美食を\n楽しんでいただける環境をつくること。\nそれがPPPJの使命です。」",
    lx + Inches(0.15), msg_y + Inches(0.05), lw - Inches(0.1), Inches(1.1),
    font_size=12, color=NAVY, italic=True, align=PP_ALIGN.LEFT)

# Right: Photo placeholder
photo_placeholder(slide, half_w + Inches(0.1), content_y, half_w - Inches(0.1), H - content_y - Inches(0.5))

slide_number(slide, 18)

# ═══════════════════════════════════════════════════
# SLIDE 19 — PPPJの違い
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "PPPJが選ばれる理由 —— 一般認証会社との違い")

content_y = Inches(1.4)

# Comparison table
tx = Inches(0.4)
tw = W - Inches(0.8)
row_h = Inches(0.62)
col1 = Inches(5.5)
col2 = Inches(2.0)
col3 = tw - col1 - col2

services = [
    ("サービス項目", "PPPJ", "一般認証会社"),
    ("ハラル認証取得サポート", "✓", "✓"),
    ("ハラル食材の安定供給", "✓", "✕"),
    ("メニュー開発・レシピ提案", "✓", "✕"),
    ("SNS・インフルエンサー集客", "✓", "✕"),
    ("Google / MEO対策", "✓", "✕"),
    ("スタッフ向けハラル教育", "✓", "✕"),
    ("継続的な改善サポート", "✓", "✕"),
]

for r, (svc, pppj, other) in enumerate(services):
    ry = content_y + r * row_h
    if r == 0:
        fill = NAVY
        fc_s = WHITE
        fc_p = GOLD
        fc_o = WHITE
    else:
        fill = LGRAY if r % 2 == 1 else WHITE
        fc_s = NAVY
        fc_p = GREEN
        fc_o = RED

    rect(slide, int(tx), int(ry), int(col1), int(row_h),
         fill_color=fill, border_color=GRAY, border_pt=0.5)
    rect(slide, int(tx + col1), int(ry), int(col2), int(row_h),
         fill_color=fill, border_color=GRAY, border_pt=0.5)
    rect(slide, int(tx + col1 + col2), int(ry), int(col3), int(row_h),
         fill_color=fill, border_color=GRAY, border_pt=0.5)

    txt(slide, svc, int(tx + Inches(0.1)), int(ry + Inches(0.12)),
        int(col1 - Inches(0.2)), int(row_h - Inches(0.15)),
        font_size=14 if r == 0 else 13, bold=(r == 0), color=fc_s, align=PP_ALIGN.LEFT)
    txt(slide, pppj, int(tx + col1), int(ry + Inches(0.12)),
        int(col2), int(row_h - Inches(0.15)),
        font_size=14 if r == 0 else 16, bold=True, color=fc_p, align=PP_ALIGN.CENTER)
    txt(slide, other, int(tx + col1 + col2), int(ry + Inches(0.12)),
        int(col3), int(row_h - Inches(0.15)),
        font_size=14 if r == 0 else 16, bold=(r == 0), color=fc_o, align=PP_ALIGN.CENTER)

slide_number(slide, 19)

# ═══════════════════════════════════════════════════
# SLIDE 20 — 6大サービス Hub
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "PPPJの6大サービス —— ワンストップで完全サポート")

content_y = Inches(1.4)

cx = W / 2
cy = content_y + (H - content_y - Inches(0.3)) / 2

# Center circle
r = Inches(1.0)
rect(slide, int(cx - r), int(cy - r * 0.75), int(r * 2), int(r * 1.5),
     fill_color=NAVY, border_color=GOLD, border_pt=3)
txt(slide, "PPPJ", cx - r, cy - r * 0.75, r * 2, r * 0.8,
    font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(slide, "ワンストップ\n支援", cx - r, cy - r * 0.75 + r * 0.8, r * 2, r * 0.65,
    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 6 surrounding service boxes
services6 = [
    ("①認証取得支援", "ハラル認証を\n確実に取得", NAVY),
    ("②食材供給", "天幸牛などの\nハラル食材調達", GREEN),
    ("③メニュー開発", "ハラルに対応した\nレシピ・メニュー提案", NAVY),
    ("④SNS集客", "インフルエンサー\n×SNS拡散", GOLD),
    ("⑤Google/MEO", "外国語検索での\n露出最大化", GREEN),
    ("⑥スタッフ教育", "現場スタッフへの\nハラル研修", NAVY),
]

positions = [
    (Inches(0.3), content_y + Inches(0.5)),   # top-left
    (Inches(0.3), content_y + Inches(2.0)),   # mid-left
    (Inches(0.3), content_y + Inches(3.5)),   # bot-left
    (W - Inches(3.0), content_y + Inches(0.5)),  # top-right
    (W - Inches(3.0), content_y + Inches(2.0)),  # mid-right
    (W - Inches(3.0), content_y + Inches(3.5)),  # bot-right
]

box_w = Inches(2.8)
box_h = Inches(1.2)

for i, ((px, py), (sname, sdesc, sc)) in enumerate(zip(positions, services6)):
    card(slide, px, py, box_w, box_h, fill=WHITE, border=sc, border_pt=2)
    rect(slide, px, py, box_w, Inches(0.1), fill_color=sc)
    txt(slide, sname, px + Inches(0.1), py + Inches(0.15), box_w - Inches(0.2), Inches(0.42),
        font_size=13, bold=True, color=sc, align=PP_ALIGN.LEFT)
    txt(slide, sdesc, px + Inches(0.1), py + Inches(0.55), box_w - Inches(0.2), Inches(0.6),
        font_size=11, color=DGRAY, align=PP_ALIGN.LEFT)
    # Arrow toward center
    if i < 3:
        arr_x = px + box_w + Inches(0.05)
        arr_y = py + box_h/2 - Inches(0.15)
        txt(slide, "→", arr_x, arr_y, Inches(0.4), Inches(0.3),
            font_size=14, color=NAVY, align=PP_ALIGN.CENTER, bold=True)
    else:
        arr_x = px - Inches(0.45)
        arr_y = py + box_h/2 - Inches(0.15)
        txt(slide, "←", arr_x, arr_y, Inches(0.4), Inches(0.3),
            font_size=14, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

slide_number(slide, 20)

# ═══════════════════════════════════════════════════
# SLIDE 21 — 天幸牛 Brand
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "天幸牛 —— PPPJが誇るハラル対応和牛ブランド")

content_y = Inches(1.4)
half_w = W / 2

# Photo placeholder left
photo_placeholder(slide, 0, content_y, half_w - Inches(0.1), H - content_y - Inches(0.5))

# Right side
rx = half_w + Inches(0.1)
rw = W - rx - Inches(0.2)

# Brand name big
txt(slide, "天幸牛",
    rx, content_y + Inches(0.2), rw, Inches(1.0),
    font_size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txt(slide, "TENKO BEEF",
    rx, content_y + Inches(1.15), rw, Inches(0.4),
    font_size=16, color=GOLD, align=PP_ALIGN.CENTER, bold=True, italic=True)

rect(slide, rx + rw*0.1, int(content_y + Inches(1.6)), int(rw * 0.8), int(Pt(3)), fill_color=GOLD)

# 4 spec boxes
specs = [
    ("産地", "国産黒毛和牛\n厳選提携牧場"),
    ("等級", "A4〜A5等級\nプレミアムグレード"),
    ("ハラル対応", "PPPJ認証処理場\n完全ハラル対応"),
    ("品質管理", "独自検査基準\nトレーサビリティ完備"),
]
spec_w = rw / 2 - Inches(0.1)
spec_h = Inches(1.2)
for j, (label, val) in enumerate(specs):
    sx = rx + (j % 2) * (spec_w + Inches(0.1))
    sy = content_y + Inches(1.75) + (j // 2) * (spec_h + Inches(0.1))
    card(slide, sx, sy, spec_w, spec_h, fill=LGRAY, border=NAVY, border_pt=1)
    txt(slide, label, sx + Inches(0.1), sy + Inches(0.1), spec_w - Inches(0.2), Inches(0.4),
        font_size=12, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    txt(slide, val, sx + Inches(0.1), sy + Inches(0.5), spec_w - Inches(0.2), Inches(0.65),
        font_size=13, color=NAVY, align=PP_ALIGN.LEFT)

slide_number(slide, 21)

# ═══════════════════════════════════════════════════
# SLIDE 22 — 天幸牛 Supply Chain
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "天幸牛 サプライチェーン —— 産地から食卓まで完全管理")

content_y = Inches(1.4)

# 5-step horizontal flow
steps5 = [
    ("提携牧場", "国産黒毛和牛\n厳選された\n提携農家"),
    ("ハラル認証\n処理場", "PPPJ提携の\n認証処理場で\n完全対応"),
    ("品質検査", "独自基準による\n品質・安全性\n確認"),
    ("PPPJ物流", "コールドチェーン\n管理での\n安全配送"),
    ("店舗提供", "認証店舗で\n提供・お客様\nへ届く"),
]
n_steps = len(steps5)
step_w = (W - Inches(0.6)) / n_steps - Inches(0.15)
step_h = Inches(3.5)
step_y = content_y + Inches(0.5)

for i, (title_s, desc) in enumerate(steps5):
    sx = Inches(0.3) + i * (step_w + Inches(0.15))
    # Main box
    box_c = [NAVY, GREEN, NAVY, GREEN, NAVY][i]
    rect(slide, int(sx), int(step_y), int(step_w), int(Inches(0.55)), fill_color=box_c)
    txt(slide, f"STEP {i+1}", sx, step_y + Inches(0.05), step_w, Inches(0.45),
        font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    card(slide, int(sx), int(step_y + Inches(0.55)), int(step_w), int(step_h - Inches(0.55)),
         fill=WHITE, border=box_c, border_pt=2)
    txt(slide, title_s, sx + Inches(0.05), step_y + Inches(0.65), step_w - Inches(0.1), Inches(0.7),
        font_size=13, bold=True, color=box_c, align=PP_ALIGN.CENTER)
    txt(slide, desc, sx + Inches(0.05), step_y + Inches(1.35), step_w - Inches(0.1), Inches(1.5),
        font_size=11, color=DGRAY, align=PP_ALIGN.CENTER)

    if i < 4:
        arr_x = sx + step_w + Inches(0.02)
        txt(slide, "→", arr_x, step_y + step_h/2 - Inches(0.2), Inches(0.12), Inches(0.4),
            font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Bottom note
rect(slide, Inches(0.3), H - Inches(0.9), W - Inches(0.6), Inches(0.6),
     fill_color=LGRAY, border_color=NAVY, border_pt=1)
txt(slide, "全工程でのトレーサビリティにより、ムスリムのお客様に「安心」を届けます",
    Inches(0.5), H - Inches(0.85), W - Inches(1.0), Inches(0.5),
    font_size=13, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

slide_number(slide, 22)

# ═══════════════════════════════════════════════════
# SLIDE 23 — 5-Step Flow
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "導入ステップ —— 5ステップで完全サポート")

content_y = Inches(1.4)

steps_flow = [
    ("STEP 1", "無料ヒアリング", "現在の状況と\nご要望を確認\n（約60分）", "〜1週間"),
    ("STEP 2", "全体設計", "最適な認証・集客\nプランを策定", "〜2週間"),
    ("STEP 3", "認証取得", "書類準備から\n認証完了まで\nサポート", "1〜3ヶ月"),
    ("STEP 4", "集客展開", "SNS・MEO・\nインフルエンサー\n施策開始", "翌月から"),
    ("STEP 5", "継続改善", "月次データ分析\n&改善提案", "継続"),
]

n = len(steps_flow)
step_w2 = (W - Inches(0.6)) / n - Inches(0.12)
step_h2 = Inches(4.2)
step_y2 = content_y + Inches(0.2)

colors5 = [NAVY, GREEN, NAVY, GREEN, NAVY]

for i, (stepnum, title_s, desc, time_e) in enumerate(steps_flow):
    sx = Inches(0.3) + i * (step_w2 + Inches(0.12))
    c = colors5[i]
    # Step number header
    rect(slide, int(sx), int(step_y2), int(step_w2), int(Inches(0.65)), fill_color=c)
    txt(slide, stepnum, sx, step_y2 + Inches(0.1), step_w2, Inches(0.5),
        font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # Card body
    card(slide, int(sx), int(step_y2 + Inches(0.65)), int(step_w2), int(step_h2 - Inches(0.65)),
         fill=WHITE, border=c, border_pt=2)
    txt(slide, title_s, sx + Inches(0.05), step_y2 + Inches(0.75), step_w2 - Inches(0.1), Inches(0.6),
        font_size=14, bold=True, color=c, align=PP_ALIGN.CENTER)
    txt(slide, desc, sx + Inches(0.05), step_y2 + Inches(1.4), step_w2 - Inches(0.1), Inches(1.6),
        font_size=12, color=DGRAY, align=PP_ALIGN.CENTER)
    # Time estimate badge
    rect(slide, int(sx + step_w2 * 0.1), int(step_y2 + step_h2 - Inches(0.7)),
         int(step_w2 * 0.8), int(Inches(0.5)), fill_color=LGRAY, border_color=c, border_pt=1)
    txt(slide, time_e, sx + step_w2 * 0.1, step_y2 + step_h2 - Inches(0.65),
        step_w2 * 0.8, Inches(0.42), font_size=11, bold=True, color=c, align=PP_ALIGN.CENTER)

    # Arrow
    if i < 4:
        arr_x = sx + step_w2 + Inches(0.01)
        txt(slide, "▶", arr_x, step_y2 + step_h2/2 - Inches(0.2), Inches(0.1), Inches(0.35),
            font_size=12, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

# Bottom note
txt(slide, "まずは無料ヒアリングから —— お気軽にご連絡ください",
    Inches(0.3), H - Inches(0.9), W - Inches(0.6), Inches(0.45),
    font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_number(slide, 23)

# ═══════════════════════════════════════════════════
# SLIDE 24 — Before/After
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "導入前 vs 導入後 —— あなたの店はどう変わるか")

content_y = Inches(1.4)
half_w = W / 2

# Left: comparison table
lx = Inches(0.3)
lw = half_w - Inches(0.4)

# Table header
row_h = Inches(0.7)
col_a = Inches(1.8)
col_b = lw - col_a * 2 - Inches(0.1)

rect(slide, int(lx), int(content_y), int(col_a), int(row_h * 0.7),
     fill_color=DGRAY, border_color=WHITE, border_pt=0.5)
rect(slide, int(lx + col_a), int(content_y), int(col_a), int(row_h * 0.7),
     fill_color=RED, border_color=WHITE, border_pt=0.5)
rect(slide, int(lx + col_a * 2), int(content_y), int(lw - col_a * 2), int(row_h * 0.7),
     fill_color=GREEN, border_color=WHITE, border_pt=0.5)

txt(slide, "項目", lx, content_y + Inches(0.08), col_a, Inches(0.5),
    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "導入前", lx + col_a, content_y + Inches(0.08), col_a, Inches(0.5),
    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "導入後（PPPJ）", lx + col_a * 2, content_y + Inches(0.08), lw - col_a * 2, Inches(0.5),
    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

ba_rows = [
    ("外国人客数", "ほぼゼロ", "月100〜330%増"),
    ("ムスリム対応", "不可（誤食リスク）", "完全対応・安心提供"),
    ("SNS・検索露出", "日本語のみ", "多言語・世界発信"),
    ("スタッフ対応", "説明できない", "研修済・自信を持って対応"),
]

for r, (item, before, after) in enumerate(ba_rows):
    ry = content_y + row_h * 0.7 + r * row_h
    fill = LGRAY if r % 2 == 0 else WHITE
    rect(slide, int(lx), int(ry), int(col_a), int(row_h),
         fill_color=fill, border_color=GRAY, border_pt=0.5)
    rect(slide, int(lx + col_a), int(ry), int(col_a), int(row_h),
         fill_color=fill, border_color=GRAY, border_pt=0.5)
    rect(slide, int(lx + col_a * 2), int(ry), int(lw - col_a * 2), int(row_h),
         fill_color=fill, border_color=GRAY, border_pt=0.5)
    txt(slide, item, lx + Inches(0.05), ry + Inches(0.1), col_a - Inches(0.1), Inches(0.5),
        font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(slide, before, lx + col_a + Inches(0.05), ry + Inches(0.1), col_a - Inches(0.1), Inches(0.5),
        font_size=11, color=RED, align=PP_ALIGN.CENTER)
    txt(slide, after, lx + col_a * 2 + Inches(0.05), ry + Inches(0.1), lw - col_a * 2 - Inches(0.1), Inches(0.5),
        font_size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Right: photo placeholder
photo_placeholder(slide, half_w + Inches(0.1), content_y, half_w - Inches(0.3), H - content_y - Inches(0.5))

slide_number(slide, 24)

# ═══════════════════════════════════════════════════
# SLIDE 25 — Closing
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, CREAM)
header_band(slide, "最後に —— 私たちが伝えたいこと")

content_y = Inches(1.4)
half_w = W / 2

# Photo placeholder left
photo_placeholder(slide, 0, content_y, half_w - Inches(0.1), H - content_y - Inches(0.4))

# Right side
rx = half_w + Inches(0.1)
rw = W - rx - Inches(0.2)

# Big message
rect(slide, rx - Inches(0.05), content_y, Inches(0.08), Inches(1.6), fill_color=GOLD)
txt(slide, "\"誰に、どう安心を届けるか\"",
    rx + Inches(0.1), content_y + Inches(0.2), rw - Inches(0.1), Inches(1.1),
    font_size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT, italic=True)

rect(slide, rx, int(content_y + Inches(1.6)), int(rw * 0.5), int(Pt(3)), fill_color=GOLD)

txt(slide, "それがすべてのビジネスの本質です。",
    rx, content_y + Inches(1.8), rw, Inches(0.5),
    font_size=14, color=NAVY, align=PP_ALIGN.LEFT)

# 3 reasons
reasons = [
    ("①", "先行者独占", "今動けば、エリアで唯一の\nハラル対応店になれる"),
    ("②", "実績証明", "神田330%など\n数字が証明している"),
    ("③", "市場拡大", "22億人のムスリム市場は\nこれからが本番"),
]
reason_y = content_y + Inches(2.4)
for j, (icon, title_r, desc) in enumerate(reasons):
    ry = reason_y + j * Inches(1.0)
    rect(slide, rx, ry, Inches(0.55), Inches(0.8), fill_color=GOLD)
    txt(slide, icon, rx, ry + Inches(0.15), Inches(0.55), Inches(0.5),
        font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(slide, title_r, rx + Inches(0.6), ry + Inches(0.05), rw - Inches(0.65), Inches(0.38),
        font_size=14, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    txt(slide, desc, rx + Inches(0.6), ry + Inches(0.42), rw - Inches(0.65), Inches(0.42),
        font_size=11, color=DGRAY, align=PP_ALIGN.LEFT)

# PPPJ message
rect(slide, rx, H - Inches(1.1), rw, Inches(0.75), fill_color=NAVY, border_color=GOLD, border_pt=1)
txt(slide, "PPPJと共に、世界に開かれた飲食店へ。",
    rx + Inches(0.1), H - Inches(1.05), rw - Inches(0.2), Inches(0.6),
    font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

slide_number(slide, 25)

# ═══════════════════════════════════════════════════
# SLIDE 26 — CTA
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg(slide, WHITE)

half_w = W / 2

# Left panel: gold background
rect(slide, 0, 0, half_w, H, fill_color=GOLD)

# Left panel content
txt(slide, "まずは、話してみてください。",
    Inches(0.4), Inches(0.6), half_w - Inches(0.5), Inches(0.9),
    font_size=20, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

rect(slide, Inches(0.4), int(Inches(1.55)), int(half_w - Inches(0.6)), int(Pt(3)), fill_color=NAVY)

consult_items = [
    "✓  ハラル認証の取り方がわからない",
    "✓  どんな食材に切り替えればいいか",
    "✓  外国人集客を始めたいがどこから",
    "✓  他店の成功事例を聞きたい",
]
for j, item in enumerate(consult_items):
    ty = Inches(1.8) + j * Inches(0.7)
    txt(slide, item, Inches(0.4), ty, half_w - Inches(0.5), Inches(0.6),
        font_size=13, color=NAVY, align=PP_ALIGN.LEFT, bold=True)

# Free badge
badge_y = Inches(4.7)
rect(slide, Inches(0.6), badge_y, half_w - Inches(1.0), Inches(0.9),
     fill_color=NAVY, border_color=WHITE, border_pt=2)
txt(slide, "初回相談  無料",
    Inches(0.6), badge_y + Inches(0.15), half_w - Inches(1.0), Inches(0.6),
    font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Bottom nav note
txt(slide, "お気軽にスタッフまでお申し付けください",
    Inches(0.4), Inches(5.8), half_w - Inches(0.5), Inches(0.4),
    font_size=12, color=NAVY, align=PP_ALIGN.LEFT, italic=True)

# Right panel: white content
rx = half_w + Inches(0.2)
rw = W - rx - Inches(0.2)

# PPPJ header
txt(slide, "PPPJ", rx, Inches(0.4), rw, Inches(0.65),
    font_size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txt(slide, "プレミアム・プロ・ピープル・ジャパン",
    rx, Inches(1.0), rw, Inches(0.4),
    font_size=13, color=DGRAY, align=PP_ALIGN.CENTER)

rect(slide, rx + rw * 0.1, int(Inches(1.45)), int(rw * 0.8), int(Pt(2)), fill_color=GOLD)

# Contact info
contact_y = Inches(1.65)
card(slide, rx, contact_y, rw, Inches(0.75), fill=LGRAY, border=NAVY)
txt(slide, "📞  047-481-8555",
    rx + Inches(0.15), contact_y + Inches(0.15), rw - Inches(0.3), Inches(0.5),
    font_size=16, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

card(slide, rx, contact_y + Inches(0.85), rw, Inches(0.75), fill=LGRAY, border=NAVY)
txt(slide, "✉  office@pppj.co.jp",
    rx + Inches(0.15), contact_y + Inches(1.0), rw - Inches(0.3), Inches(0.5),
    font_size=15, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# QR placeholder
qr_y = contact_y + Inches(1.75)
qr_size = Inches(2.0)
qr_x = rx + rw/2 - qr_size/2
photo_placeholder(slide, qr_x, qr_y, qr_size, qr_size, "📷 QRコード")

# Today's participant note
note_y = qr_y + qr_size + Inches(0.2)
rect(slide, rx, note_y, rw, Inches(0.75), fill_color=NAVY, border_color=GOLD, border_pt=2)
txt(slide, "今日の参加者限定：個別相談 優先枠あり",
    rx + Inches(0.1), note_y + Inches(0.15), rw - Inches(0.2), Inches(0.5),
    font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

slide_number(slide, 26)

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
output_path = "/home/user/20260511_ccctest/PPPJ_ハラルセミナー資料_v2.pptx"
prs.save(output_path)
print(f"✓ Saved: {output_path}")
print(f"  Slides: {len(prs.slides)}")
